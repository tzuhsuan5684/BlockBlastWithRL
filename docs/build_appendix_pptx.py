# -*- coding: utf-8 -*-
"""Build appendix_env_diagrams.pptx — two native, editable 16:9 slides that
mirror the project's deck style (Office accent1 banner, blue header table,
Calibri). Slide A: 40-D heuristic vector. Slide B: shared coordinate system.

Run:  uv run --with python-pptx python docs/build_appendix_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

EMU_PER_INCH = 914400
SCALE = 13.333 / 1280.0                      # inches per layout-px (16:9)
def X(px): return Emu(int(px * SCALE * EMU_PER_INCH))


def _set_dash(shape, val='dash'):
    ln = shape.line._get_or_add_ln()
    for el in ln.findall(qn('a:prstDash')):
        ln.remove(el)
    ln.append(ln.makeelement(qn('a:prstDash'), {'val': val}))

# ---- palette (matches the HTML / the deck) ----
BLUE      = RGBColor(0x44, 0x72, 0xC4)
BLUE_DARK = RGBColor(0x2E, 0x55, 0x97)
SEG = [RGBColor(0x2F,0x55,0x97), RGBColor(0x3D,0x63,0xB0), RGBColor(0x44,0x72,0xC4),
       RGBColor(0x5E,0x84,0xD0), RGBColor(0x8F,0xAA,0xDC), RGBColor(0xB4,0xC7,0xE7)]
ROW_ALT = RGBColor(0xD9,0xE1,0xF2)
ROW_ODD = RGBColor(0xF4,0xF6,0xFB)
ORANGE  = RGBColor(0xED,0x7D,0x31)
INK     = RGBColor(0x40,0x40,0x40)
GREY    = RGBColor(0xA6,0xA6,0xA6)
GREY2   = RGBColor(0x6A,0x6A,0x6A)
GREY3   = RGBColor(0x5A,0x5A,0x5A)
WHITE   = RGBColor(0xFF,0xFF,0xFF)
DARK    = RGBColor(0x1C,0x1C,0x1C)
CELL_BG = RGBColor(0xFA,0xFB,0xFD)
CELL_LN = RGBColor(0xC9,0xCF,0xDB)
CODEBG  = RGBColor(0xEE,0xF1,0xF8)
CODEFG  = RGBColor(0x2B,0x3A,0x55)
BOXLN   = RGBColor(0xCD,0xD6,0xEA)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def rect(slide, x, y, w, h, fill=None, line=None, lw=1.0, dash=None,
         shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, X(x), X(y), X(w), X(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
        if dash:
            _set_dash(sp)
    return sp


def rich(holder, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, m=0):
    """paras: list of paragraphs; each is list of (text,size,color,bold,font)."""
    tf = holder.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for s in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{s}", Pt(m))
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for (t, size, color, bold, font) in runs:
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.name = font
            r.font.bold = bold; r.font.color.rgb = color
    return holder


def tb(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, m=0):
    return rich(slide.shapes.add_textbox(X(x), X(y), X(w), X(h)),
                paras, align, anchor, m)


def R(t, size, color, bold=False, font='Calibri'):   # run shorthand
    return (t, size, color, bold, font)


def banner(slide, subtitle, pageno):
    bn = rect(slide, 44, 38, 1192, 70, fill=BLUE)
    rich(bn, [[R("Appendix", 24, WHITE, True)]], PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    bn.text_frame.margin_left = Pt(14)
    tb(slide, 48, 126, 1140, 42, [[R(subtitle, 19, BLUE_DARK, True)]],
       anchor=MSO_ANCHOR.MIDDLE)
    tb(slide, 1180, 690, 70, 24, [[R(pageno, 11, GREY)]], PP_ALIGN.RIGHT)


def grid(slide, ox, oy, rows, cols, cell, lab, filled, anchor):
    # column headers
    for c in range(cols):
        tb(slide, ox + lab + c*cell, oy, cell, lab,
           [[R(str(c), 10, GREY, True)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    # row headers
    for r in range(rows):
        tb(slide, ox, oy + lab + r*cell, lab, cell,
           [[R(str(r), 10, GREY, True)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    # cells
    for r in range(rows):
        for c in range(cols):
            cx, cy = ox + lab + c*cell, oy + lab + r*cell
            key = (r, c)
            if key in filled:
                cl = rect(slide, cx, cy, cell, cell, fill=BLUE, line=CELL_LN, lw=0.75)
                rich(cl, [[R(str(filled[key]), 13, WHITE, True)]],
                     PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
            else:
                rect(slide, cx, cy, cell, cell, fill=CELL_BG, line=CELL_LN, lw=0.75)
            if key == anchor:
                rect(slide, cx+2, cy+2, cell-4, cell-4, fill=None,
                     line=ORANGE, lw=2.0, dash=True)


# =====================================================================
# SLIDE A — 40-D heuristic vector
# =====================================================================
sa = prs.slides.add_slide(BLANK)
banner(sa, "Observation — 40-D Heuristic Feature Vector", "A")

tb(sa, 48, 176, 1184, 24,
   [[R("Six blocks concatenated into one length-40 vector, each normalized to "
       "[0,1], all computed on the ", 12, GREY3),
     R("current", 12, BLUE_DARK, True), R(" board (state):", 12, GREY3)]])

# segmented bar (6 equal segments, like the flex row)
segs = [("[0:8]",   ["column", "heights"],  "8 dims · ÷8"),
        ("[8:16]",  ["holes", "per column"],"8 dims · ÷7"),
        ("[16:24]", ["row", "fill counts"], "8 dims · ÷8"),
        ("[24:32]", ["column", "fill counts"],"8 dims · ÷8"),
        ("[32:39]", ["bumpiness"],          "7 dims · ÷8"),
        ("[39:40]", ["n_legal"],            "1 dim · ÷192")]
bx, bw = 48, 1184
seg_w = bw / 6.0
for i, (rng, nm, nz) in enumerate(segs):
    fg = DARK if i == 5 else WHITE
    sp = rect(sa, bx + i*seg_w, 210, seg_w - 2, 96, fill=SEG[i])
    paras = [[R(rng, 10, fg)]] + [[R(n, 12, fg, True)] for n in nm] + [[R(nz, 10, fg)]]
    rich(sp, paras, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

tb(sa, 48, 310, 1184, 20,
   [[R("← index 0  " + "·"*60 + "  index 39 →", 10, GREY)]],
   PP_ALIGN.CENTER)

# table
rows_data = [
    ("column heights", "8", "Top filled cell → bottom, per column (no gravity: counts gaps)", "_column_heights()"),
    ("holes per column", "8", "Empty cells underneath a filled cell, per column", "_holes_per_col()"),
    ("row fill counts", "8", "Filled cells in each of the 8 rows", "board.sum(axis=1)"),
    ("column fill counts", "8", "Filled cells in each of the 8 columns", "board.sum(axis=0)"),
    ("bumpiness", "7", "|height diff| of each adjacent column pair (kept per-pair, not summed)", "abs(diff(heights))"),
    ("n_legal", "1", "Number of currently legal actions — a “danger” signal", "count_nonzero(mask)"),
]
headers = ["Block", "Dims", "Meaning", "Source"]
colw = [150, 70, 734, 230]
gf = sa.shapes.add_table(len(rows_data)+1, 4, X(48), X(342), X(sum(colw)), X(210))
tbl = gf.table
tbl.first_row = False
tbl.horz_banding = False
for j, w in enumerate(colw):
    tbl.columns[j].width = X(w)

def cell_set(cell, text, size, color, bold, fill, font='Calibri', align=PP_ALIGN.LEFT):
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(7); tf.margin_right = Pt(5)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font

for j, htxt in enumerate(headers):
    cell_set(tbl.cell(0, j), htxt, 12, WHITE, True, BLUE)
for ri, row in enumerate(rows_data, start=1):
    bg = ROW_ALT if ri % 2 == 0 else ROW_ODD
    cell_set(tbl.cell(ri, 0), row[0], 12, INK, False, bg)
    cell_set(tbl.cell(ri, 1), row[1], 12, INK, False, bg, align=PP_ALIGN.CENTER)
    cell_set(tbl.cell(ri, 2), row[2], 12, INK, False, bg)
    cell_set(tbl.cell(ri, 3), row[3], 10.5, CODEFG, False, bg, font='Consolas')

tb(sa, 48, 582, 1184, 28,
   [[R("Note: ", 12, BLUE_DARK, True),
     R("distinct from the 9 afterstate features — these are per-row / "
       "per-column, ", 12, GREY3),
     R("normalized", 12, BLUE_DARK, True),
     R(", on the ", 12, GREY3), R("current", 12, BLUE_DARK, True),
     R(" state, and feed the ", 12, GREY3),
     R("critic", 12, BLUE_DARK, True), R(" V(s).", 12, GREY3)]])

# =====================================================================
# SLIDE B — shared coordinate system
# =====================================================================
sb = prs.slides.add_slide(BLANK)
banner(sb, "Piece Rendering ↔ Placement: One Shared Anchor", "B")

tb(sb, 48, 176, 1184, 24,
   [[R("The same offsets ", 12, GREY3), R("(dr, dc)", 12, CODEFG, True, 'Consolas'),
     R(" define both what the agent sees and where the piece lands.  Anchor = "
       "bounding-box ", 12, GREY3), R("top-left", 12, BLUE_DARK, True),
     R(" = cell ", 12, GREY3), R("(0,0)", 12, CODEFG, True, 'Consolas'),
     R("   (example piece: ", 12, GREY3), R("J-0", 12, BLUE_DARK, True),
     R(").", 12, GREY3)]])

# left panel — 5x5
rich(sb.shapes.add_textbox(X(55), X(244), X(252), X(46)),
     [[R("What the agent SEES", 14, BLUE_DARK, True)],
      [R("pieces[i] — 5×5 grid", 10.5, GREY2)]],
     PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
grid(sb, 55, 300, 5, 5, 46, 22, {(0,1):1,(1,1):2,(2,0):3,(2,1):4}, (0,0))
tb(sb, 50, 558, 320, 24,
   [[R("shape_to_grid(): grid[dr][dc] = 1", 10.5, CODEFG, False, 'Consolas')]],
   PP_ALIGN.CENTER)

# middle — mapping
fbox = rect(sb, 500, 312, 236, 96, fill=CODEBG, line=BOXLN, lw=1.0,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rich(fbox, [
    [R("action (", 13, CODEFG, False, 'Consolas'),
     R("row=2, col=3", 13, ORANGE, True, 'Consolas'),
     R(")", 13, CODEFG, False, 'Consolas')],
    [R(" ", 7, CODEFG, False, 'Consolas')],
    [R("board[", 12.5, CODEFG, False, 'Consolas'),
     R("row", 12.5, ORANGE, True, 'Consolas'),
     R("+dr][", 12.5, CODEFG, False, 'Consolas'),
     R("col", 12.5, ORANGE, True, 'Consolas'),
     R("+dc]=1", 12.5, CODEFG, False, 'Consolas')],
], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
rect(sb, 560, 426, 96, 48, fill=BLUE, shape=MSO_SHAPE.RIGHT_ARROW)
tb(sb, 500, 492, 236, 44,
   [[R("same (dr,dc)", 10.5, GREY2)], [R("just shifted by the anchor", 10.5, GREY2)]],
   PP_ALIGN.CENTER)

# right panel — 8x8
rich(sb.shapes.add_textbox(X(895), X(208), X(342), X(46)),
     [[R("Where it is PLACED", 14, BLUE_DARK, True)],
      [R("board — 8×8 grid", 10.5, GREY2)]],
     PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
grid(sb, 895, 272, 8, 8, 40, 22, {(2,4):1,(3,4):2,(4,3):3,(4,4):4}, (2,3))
tb(sb, 895, 616, 342, 24,
   [[R("step(): writes (row+dr, col+dc)", 10.5, CODEFG, False, 'Consolas')]],
   PP_ALIGN.CENTER)

# legend + why-it-matters
rect(sb, 48, 640, 14, 14, fill=BLUE)
tb(sb, 68, 634, 360, 24,
   [[R("piece cell  (number = same cell on both grids)", 10.5, GREY2)]],
   anchor=MSO_ANCHOR.MIDDLE)
rect(sb, 470, 640, 14, 14, fill=None, line=ORANGE, lw=2.0,
     dash=True)
tb(sb, 490, 634, 760, 24,
   [[R("anchor (0,0) / (row,col) — ", 10.5, GREY2),
     R("may be empty", 10.5, BLUE_DARK, True),
     R("; padding always sits bottom-right", 10.5, GREY2)]],
   anchor=MSO_ANCHOR.MIDDLE)
tb(sb, 48, 670, 1184, 26,
   [[R("Why it matters: ", 12, BLUE_DARK, True),
     R("the relative shape the agent sees in the 5×5 patch equals the "
       "relative footprint it places on the board — no hidden coordinate "
       "conversion to learn.", 12, GREY3)]])

out = "docs/appendix_env_diagrams.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
