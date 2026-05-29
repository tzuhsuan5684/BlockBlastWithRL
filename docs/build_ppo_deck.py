#!/usr/bin/env python
"""Build the PPO section deck (Section 05) in the Block Blast RL house style.

Reproduces the visual language of docs/Presentation.pdf:
  - deep navy background on the divider, light blue-grey on content slides
  - teal title band with a "PPO · SECTION" kicker, page number, footer rule
  - white cards with a coloured left accent edge; dark navy "hero" cards
  - amber as the single highlight colour

Outputs (next to this script):
  - PPO_Section05.pptx   English on slides, Chinese in speaker notes
  - PPO_speaker_script.docx   full EN + 中文 narration, one block per slide

Run (no project deps touched):
    uv run --with python-pptx --with python-docx python docs/build_ppo_deck.py
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))

# ----------------------------------------------------------------------------
# Palette (sampled from Presentation.pdf)
# ----------------------------------------------------------------------------
NAVY    = RGBColor(0x0F, 0x3A, 0x4D)   # divider bg / hero cards
NAVY_CIRCLE = RGBColor(0x15, 0x55, 0x6B)
BAR     = RGBColor(0x1C, 0x79, 0x82)   # title band
TEAL    = RGBColor(0x17, 0xA2, 0x94)   # bright accent / numbers / left edges
TEAL_DK = RGBColor(0x0E, 0x7C, 0x86)
TEAL_LT = RGBColor(0x5F, 0xB0, 0xAB)   # subtitle teal on navy
AMBER   = RGBColor(0xE8, 0xA2, 0x3C)   # single highlight
LIGHTBG = RGBColor(0xED, 0xF2, 0xF5)   # content slide bg
CARD    = RGBColor(0xFF, 0xFF, 0xFF)
INK     = RGBColor(0x2E, 0x3A, 0x42)   # body text
MUTE    = RGBColor(0x7C, 0x8A, 0x94)   # secondary text
BORDER  = RGBColor(0xDC, 0xE5, 0xEA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
INK_ON_DARK = RGBColor(0xCF, 0xDD, 0xE2)

LATIN = "Lato"          # PowerPoint substitutes gracefully if absent
MONO  = "Consolas"

FONT_SCALE = 1.13       # global "make everything a bit bigger" knob

# ----------------------------------------------------------------------------
# Geometry helpers — the reference canvas is 1333x750 px == 13.33x7.5 in,
# so pixel/100 == inches and we can read coordinates straight off the slides.
# ----------------------------------------------------------------------------
def px(v: float) -> Emu:
    return Inches(v / 100.0)

PAGE_W, PAGE_H = 1333, 750


def _set_fill(shape, color):
    if color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color


def _set_line(shape, color, width=1.0):
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color
        shape.line.width = Pt(width)


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, rounded=False,
         radius=0.045, shadow=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        px(x), px(y), px(w), px(h))
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    _set_fill(shp, fill)
    _set_line(shp, line, line_w)
    shp.shadow.inherit = False
    if shadow:
        _soft_shadow(shp)
    return shp


def oval(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(x), px(y), px(w), px(h))
    _set_fill(shp, fill)
    _set_line(shp, None)
    shp.shadow.inherit = False
    return shp


def _soft_shadow(shape):
    # subtle drop shadow on white cards
    spPr = shape._element.spPr
    effLst = spPr.makeelement(qn('a:effectLst'), {})
    sh = spPr.makeelement(qn('a:outerShdw'),
                          {'blurRad': '90000', 'dist': '38000',
                           'dir': '5400000', 'rotWithShape': '0'})
    clr = spPr.makeelement(qn('a:srgbClr'), {'val': '6E7F8A'})
    alpha = spPr.makeelement(qn('a:alpha'), {'val': '22000'})
    clr.append(alpha)
    sh.append(clr)
    effLst.append(sh)
    spPr.append(effLst)


def txt(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    """paras: list of paragraph dicts.

    Each para: {runs:[{t,size,bold,italic,color,font,spacing}], align, space_before,
                space_after, line}  OR a shorthand single-run dict {t,size,...}.
    """
    box = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, ):
        pass
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        if "space_after" in para:
            p.space_after = Pt(para["space_after"])
        if "space_before" in para:
            p.space_before = Pt(para["space_before"])
        if "line" in para:
            p.line_spacing = para["line"]
        runs = para.get("runs")
        if runs is None:
            runs = [para]
        for rd in runs:
            r = p.add_run()
            r.text = rd.get("t", "")
            f = r.font
            f.size = Pt(rd.get("size", 14) * FONT_SCALE)
            f.bold = rd.get("bold", False)
            f.italic = rd.get("italic", False)
            f.name = rd.get("font", LATIN)
            f.color.rgb = rd.get("color", INK)
            if "spacing" in rd:
                _letter_spacing(r, rd["spacing"])
    return box


def _letter_spacing(run, pts):
    run.font._rPr.set('spc', str(int(pts * 100)))


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ----------------------------------------------------------------------------
# Reusable slide furniture
# ----------------------------------------------------------------------------
def content_bg(slide):
    rect(slide, 0, 0, PAGE_W, PAGE_H, fill=LIGHTBG)


def title_band(slide, kicker, title, pageno):
    # teal band
    rect(slide, 0, 0, PAGE_W, 112, fill=BAR)
    rect(slide, 0, 0, 14, 112, fill=TEAL)            # brighter left edge
    txt(slide, 60, 20, 900, 24,
        [{"t": kicker, "size": 11.5, "bold": True, "color": TEAL_LT,
          "font": LATIN, "spacing": 1.4}])
    txt(slide, 60, 40, 1120, 60,
        [{"t": title, "size": 26, "bold": True, "color": WHITE, "font": LATIN}],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, PAGE_W - 200, 40, 150, 28,
        [{"t": pageno, "size": 12, "color": RGBColor(0xBF, 0xD7, 0xDB),
          "align": PP_ALIGN.RIGHT}], anchor=MSO_ANCHOR.MIDDLE)
    # footer
    rect(slide, 60, 706, PAGE_W - 120, 1, fill=BORDER)
    txt(slide, 60, 712, 700, 22,
        [{"t": "Block Blast RL  ·  PPO (Member C)", "size": 10.5, "color": MUTE,
          "font": LATIN}])


def left_accent(slide, x, y, h, color, w=6):
    rect(slide, x, y, w, h, fill=color)


def hero_card(slide, x, y, w, h, kicker, big, sub=None, body=None,
              accent=AMBER):
    rect(slide, x, y, w, h, fill=NAVY, rounded=True, radius=0.04)
    cy = y + 32
    txt(slide, x + 34, cy, w - 68, 18,
        [{"t": kicker, "size": 11, "bold": True, "color": TEAL_LT,
          "spacing": 1.2}])
    cy += 26
    txt(slide, x + 34, cy, w - 68, 104,
        [{"t": big, "size": 27, "bold": True, "color": WHITE, "line": 1.04}])
    cy += 112
    rect(slide, x + 36, cy, 48, 5, fill=accent)
    cy += 20
    if sub:
        txt(slide, x + 34, cy, w - 68, 46,
            [{"t": sub, "size": 13.5, "color": TEAL_LT, "line": 1.2}])
        cy += 52
    if body:
        txt(slide, x + 34, cy, w - 68, y + h - cy - 24, body)


# ----------------------------------------------------------------------------
# Deck content
# ----------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(BLANK)


# ===== Slide 0 — Section divider =============================================
s = new_slide()
rect(s, 0, 0, PAGE_W, PAGE_H, fill=NAVY)
oval(s, 905, 350, 540, 540, NAVY_CIRCLE)
oval(s, -190, 470, 430, 430, NAVY_CIRCLE)
oval(s, 815, 44, 26, 26, AMBER)
txt(s, 95, 205, 600, 24,
    [{"t": "FINAL REPORT  ·  SECTION 05", "size": 13, "bold": True,
      "color": TEAL, "spacing": 2.0}])
txt(s, 90, 235, 1100, 120,
    [{"t": "05 · PPO", "size": 74, "bold": True, "color": WHITE}])
txt(s, 95, 372, 1050, 44,
    [{"t": "Proximal Policy Optimization", "size": 27, "color": TEAL_LT}])
rect(s, 98, 432, 70, 6, fill=AMBER)
txt(s, 95, 460, 1100, 30,
    [{"t": "Algorithm · the math that makes it work · a from-scratch "
           "implementation for Block Blast", "size": 15, "color": INK_ON_DARK}])
txt(s, 95, 500, 1100, 26,
    [{"t": "Member C  ·  PPO", "size": 14, "color": TEAL_LT}])
txt(s, 95, 700, 700, 24,
    [{"t": "National Central University  ·  RL Final Project",
      "size": 11.5, "color": RGBColor(0x6A, 0x86, 0x92)}])
notes(s, "章節頁。這個章節講 PPO——Proximal Policy Optimization，也就是我為我們 "
         "agent 用的演算法。我會講它是哪一類方法、讓它運作的數學、以及我怎麼從零實作它。")

# ===== Slide 1 — What kind of algorithm =====================================
s = new_slide()
content_bg(s)
title_band(s, "PPO · POSITIONING", "What Kind of Algorithm Is PPO?", "01 / 14")
cols = [
    ("Model-Free", "Learns purely from experience — never models the game's "
     "transitions.", TEAL),
    ("Policy-Based", "Actor–critic: directly optimizes the policy π(a|s), "
     "not Q-values.", TEAL),
    ("On-Policy", "Trains on data the current policy just collected, then "
     "discards it.", AMBER),
]
cw, gap, x0, cy = 388, 26, 60, 165
for i, (h, b, ac) in enumerate(cols):
    cx = x0 + i * (cw + gap)
    rect(s, cx, cy, cw, 60, fill=NAVY if ac is TEAL else NAVY, rounded=False)
    rect(s, cx, cy, cw, 60, fill=ac)
    txt(s, cx + 24, cy, cw - 40, 60,
        [{"t": h, "size": 20, "bold": True, "color": WHITE}],
        anchor=MSO_ANCHOR.MIDDLE)
    rect(s, cx, cy + 60, cw, 168, fill=CARD, shadow=True)
    txt(s, cx + 24, cy + 88, cw - 48, 120,
        [{"t": b, "size": 14.5, "color": INK, "line": 1.25}])
# bottom banner
by = 470
rect(s, 60, by, PAGE_W - 120, 150, fill=NAVY, rounded=True, radius=0.03)
txt(s, 92, by + 26, PAGE_W - 200, 20,
    [{"t": "SAME FAMILY AS A2C / A3C", "size": 11, "bold": True,
      "color": TEAL_LT, "spacing": 1.4}])
txt(s, 92, by + 52, PAGE_W - 184, 80,
    [{"runs": [
        {"t": "PPO is the same actor–critic family as A2C/A3C — ", "size": 16,
         "color": WHITE},
        {"t": "but noticeably more stable and more sample-efficient.",
         "size": 16, "bold": True, "color": WHITE}], "line": 1.25},
     {"runs": [
        {"t": "Contrast with our DQN — ", "size": 15, "color": INK_ON_DARK},
        {"t": "value-based + off-policy", "size": 15, "bold": True,
         "color": AMBER},
        {"t": ".  That contrast is exactly what the project compares.",
         "size": 15, "color": INK_ON_DARK}], "space_before": 8, "line": 1.25}])
notes(s, "PPO 的定位。它是 model-free——純粹從經驗學，從不預測遊戲動態。它是 "
         "policy-based，直接優化策略（actor），而不是像 DQN 那樣學 Q 值。它是 "
         "on-policy——用當前策略剛收集的資料訓練，然後丟掉。和 A2C/A3C 同家族，"
         "但更穩定、樣本效率更高。對比 DQN 的 value-based + off-policy，這正是"
         "我們專案要比較的兩端。")

# ===== Slide 2 — Training pipeline ==========================================
s = new_slide()
content_bg(s)
title_band(s, "PPO · TRAINING LOOP", "The PPO Training Pipeline", "02 / 14")
# left hero with the 3 stages
hero_card(s, 60, 150, 560, 470,
          "ONE ITERATION = THREE STAGES", "Collect → GAE →\nOptimize",
          sub="Repeat until trained · network frozen during GAE",
          body=[
              {"runs": [{"t": "1  ROLLOUT", "size": 13.5, "bold": True,
                         "color": AMBER}], "space_after": 2},
              {"t": "actor plays N steps; store s, a, logπ_old, r, V, done",
               "size": 12.5, "color": INK_ON_DARK, "line": 1.15,
               "space_after": 13},
              {"runs": [{"t": "2  GAE", "size": 13.5, "bold": True,
                         "color": AMBER}], "space_after": 2},
              {"t": "backward pass → advantages Â and returns R̂",
               "size": 12.5, "color": INK_ON_DARK, "line": 1.15,
               "space_after": 13},
              {"runs": [{"t": "3  OPTIMIZE", "size": 13.5, "bold": True,
                         "color": AMBER}], "space_after": 2},
              {"t": "shuffle; K epochs of minibatch updates on the SAME batch",
               "size": 12.5, "color": INK_ON_DARK, "line": 1.15},
          ])
# right: efficiency + the hook
rx = 652
rect(s, rx, 150, 621, 150, fill=CARD, shadow=True)
left_accent(s, rx, 150, 150, TEAL)
txt(s, rx + 28, 172, 560, 22,
    [{"t": "vs A2C — the efficiency win", "size": 16, "bold": True,
      "color": TEAL_DK}])
txt(s, rx + 28, 202, 565, 90,
    [{"t": "Stage 3 reuses the same batch for K epochs instead of discarding "
           "it after a single update — far more sample-efficient.",
      "size": 14, "color": INK, "line": 1.3}])
rect(s, rx, 318, 621, 302, fill=NAVY, rounded=True, radius=0.03)
txt(s, rx + 28, 344, 565, 22,
    [{"t": "BUT — IS THAT SAFE?", "size": 11, "bold": True, "color": AMBER,
      "spacing": 1.4}])
txt(s, rx + 28, 372, 565, 70,
    [{"t": "Reusing on-policy data is exactly what breaks vanilla policy "
           "gradient. So how does PPO get away with it?", "size": 15.5,
      "color": WHITE, "line": 1.3}])
txt(s, rx + 28, 452, 565, 22,
    [{"t": "TWO TRICKS — COMING UP IN THE MATH", "size": 11, "bold": True,
      "color": TEAL_LT, "spacing": 1.2}])
txt(s, rx + 28, 480, 565, 120,
    [{"runs": [{"t": "Importance sampling (rₜ)", "size": 14.5, "bold": True,
                "color": WHITE},
               {"t": "  →  makes the reuse correct", "size": 14.5,
                "color": INK_ON_DARK}], "space_after": 8},
     {"runs": [{"t": "Clipping", "size": 14.5, "bold": True, "color": WHITE},
               {"t": "  →  keeps every reuse step small → stable",
                "size": 14.5, "color": INK_ON_DARK}]}])
notes(s, "進數學前先看 PPO 怎麼訓練——三階段迴圈。階段一 rollout：actor 跑 N 步，"
         "存狀態、動作、舊 log 機率、獎勵、價值、done。階段二 GAE：凍結網路，一次"
         "反向掃描算出優勢與回報。階段三 optimize：打散，用同一批資料做 K 個 epoch。"
         "重點：第三階段把同一批資料重複利用 K 個 epoch，而非更新一次就丟——效率大增，"
         "但重複利用 on-policy 資料正是搞垮傳統 policy gradient 的元兇。那 PPO 怎麼做到？"
         "兩招——重要性採樣與 clip——正是接下來數學要講的。")

# ===== Slide 3 — Math at a glance ===========================================
s = new_slide()
content_bg(s)
title_band(s, "PPO · MATH OVERVIEW", "The Math at a Glance", "03 / 14")
txt(s, 60, 130, 1180, 26,
    [{"t": "Two tricks plus two more pieces — four formulas. We walk through "
           "each next.", "size": 14.5, "color": MUTE}])
forms = [
    ("1  Probability ratio", "how much the policy changed",
     "rₜ(θ) = π_new(aₜ|sₜ) / π_old(aₜ|sₜ)", TEAL),
    ("2  Advantage (GAE)", "was the action better than expected?",
     "Âₜ = δₜ + γλ·Âₜ₊₁,   δₜ = rₜ + γV(sₜ₊₁) − V(sₜ)", TEAL),
    ("3  Clipped surrogate", "the safe policy update",
     "L^CLIP = Ê[ min( rₜÂₜ , clip(rₜ,1−ε,1+ε)Âₜ ) ]", AMBER),
    ("4  Total objective", "actor + critic + exploration  (maximize)",
     "L^obj = L^CLIP − c₁·L^VF + c₂·S", TEAL),
]
cw, ch, gx, gy, x0, y0 = 588, 192, 24, 22, 60, 178
for i, (h, sub, f, ac) in enumerate(forms):
    cx = x0 + (i % 2) * (cw + gx)
    cyy = y0 + (i // 2) * (ch + gy)
    rect(s, cx, cyy, cw, ch, fill=CARD, shadow=True)
    left_accent(s, cx, cyy, ch, ac)
    txt(s, cx + 28, cyy + 24, cw - 50, 24,
        [{"t": h, "size": 18, "bold": True, "color": INK}])
    txt(s, cx + 28, cyy + 54, cw - 50, 22,
        [{"t": sub, "size": 13, "italic": True, "color": MUTE}])
    rect(s, cx + 28, cyy + 92, cw - 56, 64, fill=LIGHTBG, rounded=True,
         radius=0.08)
    txt(s, cx + 40, cyy + 92, cw - 80, 64,
        [{"t": f, "size": 15.5, "bold": True, "color": TEAL_DK,
          "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
notes(s, "我們剛看完迴圈和兩招，這頁把整個數學放一起。PPO 其實就四塊：量化策略改變"
         "多少的 ratio、量化動作好不好的 advantage、安全更新策略的裁剪目標、把 "
         "critic 和探索綁進來、要最大化的總目標。符號先別擔心，一個一個講。這頁當地圖。")

# ===== Slide 4 — Piece 1 ratio ==============================================
s = new_slide()
content_bg(s)
title_band(s, "PPO · MATH 1 / 4", "Probability Ratio = Importance Sampling",
           "04 / 14")
hero_card(s, 60, 150, 560, 470,
          "PIECE 1 · THE RATIO", "rₜ(θ)",
          sub="new-policy prob ÷ old-policy prob, same action",
          accent=TEAL,
          body=[
              {"t": "rₜ(θ) = π_new(aₜ|sₜ) / π_old(aₜ|sₜ)", "size": 16,
               "bold": True, "color": WHITE, "line": 1.2, "space_after": 18},
              {"runs": [{"t": "rₜ = 1", "size": 14, "bold": True,
                         "color": AMBER},
                        {"t": "    unchanged", "size": 14,
                         "color": INK_ON_DARK}], "space_after": 10},
              {"runs": [{"t": "rₜ > 1", "size": 14, "bold": True,
                         "color": AMBER},
                        {"t": "    new policy prefers it more", "size": 14,
                         "color": INK_ON_DARK}], "space_after": 10},
              {"runs": [{"t": "rₜ < 1", "size": 14, "bold": True,
                         "color": AMBER},
                        {"t": "    prefers it less", "size": 14,
                         "color": INK_ON_DARK}]},
          ])
cards = [
    ("It corrects for reused data", "We're reusing data the OLD policy "
     "collected; this ratio is what makes that reuse mathematically valid — "
     "the first of the two tricks.", TEAL),
    ("Computed in log-space", "Subtract log-probs, then exponentiate — "
     "numerically stable because raw probabilities can be tiny.", TEAL),
    ("It is also what gets clipped", "PPO will clip this very quantity — the "
     "second trick, two slides from here.", AMBER),
]
rx, ry = 652, 150
for h, b, ac in cards:
    hh = 150
    rect(s, rx, ry, 621, hh, fill=CARD, shadow=True)
    left_accent(s, rx, ry, hh, ac)
    txt(s, rx + 28, ry + 22, 565, 24,
        [{"t": h, "size": 16, "bold": True, "color": INK}])
    txt(s, rx + 28, ry + 54, 568, 88,
        [{"t": b, "size": 13.5, "color": INK, "line": 1.3}])
    ry += hh + 10
notes(s, "第一塊 ratio，也是第一招重要性採樣。新策略對某動作的機率，除以舊策略對"
         "同動作的機率。等於 1 沒變、大於 1 更喜歡。因為在重複利用舊策略的資料，這個"
         "比率正是修正落差、讓重複利用在數學上成立的東西。用 log 空間算以保數值穩定。"
         "這個 ratio 也是接下來要 clip 的量——第二招。")

# ===== Slide 5 — Piece 2 advantage / GAE ====================================
s = new_slide()
content_bg(s)
title_band(s, "PPO · MATH 2 / 4", "Advantage & GAE", "05 / 14")
hero_card(s, 60, 150, 560, 470,
          "PIECE 2 · ADVANTAGE", "Was it better\nthan expected?",
          sub="how much better than the critic's baseline V(s)",
          accent=TEAL,
          body=[
              {"runs": [{"t": "Âₜ > 0", "size": 14, "bold": True,
                         "color": AMBER},
                        {"t": "  reinforce         ", "size": 14,
                         "color": INK_ON_DARK},
                        {"t": "Âₜ < 0", "size": 14, "bold": True,
                         "color": AMBER},
                        {"t": "  discourage", "size": 14,
                         "color": INK_ON_DARK}], "space_after": 18},
              {"t": "δₜ = rₜ + γV(sₜ₊₁) − V(sₜ)", "size": 14, "bold": True,
               "color": WHITE, "space_after": 9},
              {"t": "Âₜ = δₜ + γλ·Âₜ₊₁    (backward in time)", "size": 14,
               "bold": True, "color": WHITE, "space_after": 9},
              {"t": "R̂ₜ = Âₜ + V(sₜ)    → critic target", "size": 13,
               "color": TEAL_LT},
          ])
rows = [
    ("one-step TD", "next real reward + critic's guess V(s′)  (bootstrap)",
     "low variance · high bias", TEAL),
    ("Monte-Carlo", "real rewards to episode end, no guessing",
     "high variance · low bias (unbiased)", TEAL),
    ("GAE-λ", "weighted blend; λ→0 ≈ TD, λ→1 ≈ MC  (we use λ=0.95)",
     "Âₜ = Σ (γλ)ˡ δₜ₊ₗ — exponentially-weighted; λ is the bias↔variance dial",
     AMBER),
]
rx, ry = 652, 150
for h, b, tag, ac in rows:
    hh = 150
    rect(s, rx, ry, 621, hh, fill=CARD, shadow=True)
    left_accent(s, rx, ry, hh, ac)
    txt(s, rx + 28, ry + 20, 565, 24,
        [{"t": h, "size": 16, "bold": True, "color": INK}])
    txt(s, rx + 28, ry + 50, 568, 44,
        [{"t": b, "size": 13, "color": INK, "line": 1.25}])
    txt(s, rx + 28, ry + 110, 568, 32,
        [{"t": tag, "size": 12, "italic": True,
          "color": (AMBER if ac is AMBER else TEAL_DK), "line": 1.15}])
    ry += hh + 10
notes(s, "第二塊 advantage——比 critic 預期好或壞？正的強化、負的抑制。怎麼估？兩個"
         "極端：單步 TD 只用下一個真實獎勵加 critic 的猜測（bootstrap），穩定、低變異"
         "但有偏；Monte-Carlo 用到結束的真實獎勵、不猜，無偏但高變異。GAE 一般化兩者，"
         "由 λ 控制（0≈TD、1≈MC，用 0.95）。把遞迴展開，advantage 就是未來 TD error "
         "的指數加權和，衰減 γλ——越遠權重越小，λ 就是變異↔偏差的旋鈕。邊界遞迴重置，"
         "價值不跨越遊戲結束。")

# ===== Slide 6 — Piece 3 clipped objective ==================================
s = new_slide()
content_bg(s)
title_band(s, "PPO · MATH 3 / 4", "The Clipped Surrogate Objective", "06 / 14")
# big formula banner
rect(s, 60, 150, PAGE_W - 120, 96, fill=NAVY, rounded=True, radius=0.03)
txt(s, 60, 150, PAGE_W - 120, 96,
    [{"t": "L^CLIP = Ê[ min( rₜÂₜ ,  clip(rₜ, 1−ε, 1+ε)·Âₜ ) ]",
      "size": 23, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}],
    anchor=MSO_ANCHOR.MIDDLE)
cards = [
    ("Unclipped  rₜÂₜ", "the ordinary policy-gradient term", TEAL),
    ("Clipped", "ratio forced into [1−ε, 1+ε] — with ε = 0.2 that is "
     "[0.8, 1.2]", TEAL),
    ("min → pessimistic", "keep the smaller term → a lower bound, so there's "
     "no incentive to move the policy too far", AMBER),
]
cw, gx, x0, cy = 388, 26, 60, 272
for i, (h, b, ac) in enumerate(cards):
    cx = x0 + i * (cw + gx)
    rect(s, cx, cy, cw, 168, fill=CARD, shadow=True)
    left_accent(s, cx, cy, 168, ac)
    txt(s, cx + 26, cy + 22, cw - 46, 26,
        [{"t": h, "size": 16.5, "bold": True, "color": INK}])
    txt(s, cx + 26, cy + 58, cw - 50, 100,
        [{"t": b, "size": 13.5, "color": INK, "line": 1.3}])
by = 470
rect(s, 60, by, PAGE_W - 120, 150, fill=NAVY, rounded=True, radius=0.03)
txt(s, 92, by + 24, PAGE_W - 184, 22,
    [{"t": "THIS IS THE SECOND TRICK — IT ANSWERS THE PIPELINE PUZZLE",
      "size": 11, "bold": True, "color": AMBER, "spacing": 1.2}])
txt(s, 92, by + 50, PAGE_W - 184, 96,
    [{"runs": [{"t": "In one line:  ", "size": 15, "color": INK_ON_DARK},
               {"t": "a ceiling on already-good moves, a floor on already-bad "
                "moves, ", "size": 15, "color": WHITE},
               {"t": "but a full corrective gradient if the policy lurches "
                "the wrong way", "size": 15, "bold": True, "color": WHITE},
               {"t": "  (each case next slide).", "size": 15,
                "color": INK_ON_DARK}], "line": 1.3},
     {"t": "A cheap stand-in for TRPO's expensive trust-region constraint.",
      "size": 13.5, "italic": True, "color": TEAL_LT, "space_before": 8}])
notes(s, "第三塊是 PPO 核心，也是第二招——流程那頁疑問的答案。同一項取兩版本：未裁剪"
         "（一般 policy gradient）與裁剪（比率被壓在 0.8 到 1.2）。取較小、較悲觀的，"
         "形成下界，移動太遠得不到獎勵——這正是讓重複利用安全的原因。一句話：好動作設"
         "天花板、壞動作設地板，但往錯方向給完整修正梯度。便宜地取代 TRPO 的 "
         "trust-region。下一頁逐一說明。")

# ===== Slide 7 — Why clipping works (cases) =================================
s = new_slide()
content_bg(s)
title_band(s, "PPO · MATH 3 / 4", "Why Clipping Works — Case Analysis",
           "07 / 14")
cases = [
    ("CASE A", "Good action  (Âₜ > 0)",
     ["rₜ ≤ 1+ε : unclipped = clipped → normal gradient, keep improving",
      "rₜ > 1+ε : min picks (1+ε)Âₜ → gradient = 0",
      "→ a CEILING — stop over-committing"], TEAL),
    ("CASE B", "Bad action  (Âₜ < 0)",
     ["rₜ ≥ 1−ε : normal gradient, keep suppressing",
      "rₜ < 1−ε : min picks (1−ε)Âₜ → gradient = 0",
      "→ a FLOOR — stop over-punishing"], TEAL),
    ("SAFETY", "Wrong way  (Âₜ < 0, rₜ jumps > 1+ε)",
     ["unclipped rₜÂₜ is MORE negative than clipped",
      "min keeps the unclipped term → large corrective gradient",
      "→ policy pulled back hard — no free pass"], AMBER),
]
cw, gx, x0, cy = 388, 26, 60, 160
for i, (k, h, lines, ac) in enumerate(cases):
    cx = x0 + i * (cw + gx)
    rect(s, cx, cy, cw, 70, fill=ac)
    txt(s, cx + 24, cy + 12, cw - 40, 18,
        [{"t": k, "size": 11, "bold": True, "color": WHITE, "spacing": 1.4}])
    txt(s, cx + 24, cy + 32, cw - 40, 30,
        [{"t": h, "size": 16.5, "bold": True, "color": WHITE}])
    rect(s, cx, cy + 70, cw, 300, fill=CARD, shadow=True)
    paras = []
    for j, ln in enumerate(lines):
        bold = ln.startswith("→")
        paras.append({"t": ln, "size": 13.5,
                      "color": (ac if bold and ac is AMBER else
                                (TEAL_DK if bold else INK)),
                      "bold": bold, "line": 1.2, "space_after": 12})
    txt(s, cx + 26, cy + 96, cw - 50, 250, paras)
by = 568
rect(s, 60, by, PAGE_W - 120, 56, fill=NAVY, rounded=True, radius=0.06)
txt(s, 60, by, PAGE_W - 120, 56,
    [{"runs": [{"t": "Summary:  ", "size": 14.5, "bold": True, "color": AMBER},
               {"t": "ceiling on good moves · floor on bad moves · no free "
                "pass for going the wrong way.", "size": 14.5,
                "color": WHITE}], "align": PP_ALIGN.CENTER}],
    anchor=MSO_ANCHOR.MIDDLE)
notes(s, "逐情況看。好動作 Â>0：比率在範圍內是正常梯度；超過 1.2，min 切常數、梯度為"
         "零——天花板，防過度押注。壞動作 Â<0 對稱：低於 0.8 梯度變平——地板，防過度"
         "懲罰。安全情況：壞動作往錯方向、比率暴衝時，未裁剪項更負，min 保留它，給大"
         "梯度把策略拉回。所以：天花板、地板，但往錯方向不給通行證。")

# ===== Slide 8 — Complete loss ==============================================
s = new_slide()
content_bg(s)
title_band(s, "PPO · MATH 4 / 4", "The Complete Loss", "08 / 14")
# objective vs loss two cards
rect(s, 60, 152, 588, 150, fill=TEAL_DK, rounded=True, radius=0.03)
txt(s, 84, 170, 540, 20,
    [{"t": "PAPER · OBJECTIVE (maximize)", "size": 11, "bold": True,
      "color": WHITE, "spacing": 1.2}])
txt(s, 84, 200, 540, 80,
    [{"t": "L^obj = L^CLIP − c₁·L^VF + c₂·S", "size": 21, "bold": True,
      "color": WHITE}], anchor=MSO_ANCHOR.MIDDLE)
rect(s, 685, 152, 588, 150, fill=AMBER, rounded=True, radius=0.03)
txt(s, 709, 170, 540, 20,
    [{"t": "CODE · LOSS (minimize)", "size": 11, "bold": True, "color": NAVY,
      "spacing": 1.2}])
txt(s, 709, 200, 540, 80,
    [{"t": "L^loss = −L^CLIP + c₁·L^VF − c₂·S", "size": 21, "bold": True,
      "color": NAVY}], anchor=MSO_ANCHOR.MIDDLE)
terms = [
    ("L^CLIP", "clipped surrogate → trains the ACTOR", TEAL),
    ("L^VF", "MSE( V(s), GAE returns ) → trains the CRITIC   ·   c₁ = 0.5",
     TEAL),
    ("S", "policy entropy → keeps exploration alive   ·   c₂ = 0.01", AMBER),
]
ty = 322
for h, b, ac in terms:
    rect(s, 60, ty, PAGE_W - 120, 64, fill=CARD, shadow=True)
    left_accent(s, 60, ty, 64, ac)
    txt(s, 86, ty, 150, 64,
        [{"t": h, "size": 17, "bold": True, "color": INK}],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 250, ty, PAGE_W - 320, 64,
        [{"t": b, "size": 14.5, "color": INK}], anchor=MSO_ANCHOR.MIDDLE)
    ty += 74
by = 556
rect(s, 60, by, PAGE_W - 120, 66, fill=NAVY, rounded=True, radius=0.05)
txt(s, 60, by, PAGE_W - 120, 66,
    [{"runs": [{"t": "Why the sign flips:  ", "size": 14, "bold": True,
                "color": AMBER},
               {"t": "the paper maximizes, PyTorch minimizes — so the code "
                "negates the whole thing. Same equation, opposite sign.",
                "size": 14, "color": WHITE}], "align": PP_ALIGN.CENTER}],
    anchor=MSO_ANCHOR.MIDDLE)
notes(s, "第四塊把全部綁起來。三項：裁剪目標訓練 actor；均方誤差訓練 critic 預測 GAE "
         "returns；熵獎勵保持探索、避免太早崩塌。常被問的細節：價值項在論文『目標』裡"
         "是減、在程式碼裡是加。那是因為論文最大化、PyTorch 最小化，程式碼把整個取負號，"
         "每項符號都翻。同一條式子、相反慣例。係數 0.5、0.01 沿用 SB3。")

# ===== Slide 9 — Action masking =============================================
s = new_slide()
content_bg(s)
title_band(s, "PPO · IMPLEMENTATION 1", "Action Masking", "09 / 14")
hero_card(s, 60, 150, 560, 470,
          "WHY IT IS MANDATORY", "Most of the 192\nare illegal",
          sub="out of bounds · cell filled · slot already used",
          accent=AMBER,
          body=[{"t": "Without masking, the policy wastes most of its capacity "
                      "just learning not to make illegal moves — and the env "
                      "may even crash.", "size": 14, "color": INK_ON_DARK,
                 "line": 1.35}])
cards = [
    ("Discrete(192)", "action = piece_idx × 64 + row × 8 + col — one flat int "
     "is the coordinate every agent speaks.", TEAL),
    ("Mask before sampling", "set illegal logits to −∞ → after softmax they "
     "get 0 probability and contribute 0 entropy.", TEAL),
    ("Game-over = no legal action", "the same mask doubles as termination: "
     "terminated = not any(mask).", AMBER),
]
rx, ry = 652, 150
for h, b, ac in cards:
    hh = 150
    rect(s, rx, ry, 621, hh, fill=CARD, shadow=True)
    left_accent(s, rx, ry, hh, ac)
    txt(s, rx + 28, ry + 22, 565, 24,
        [{"t": h, "size": 16, "bold": True, "color": INK}])
    txt(s, rx + 28, ry + 54, 568, 88,
        [{"t": b, "size": 13.5, "color": INK, "line": 1.3}])
    ry += hh + 10
notes(s, "實作最重要的細節是動作遮罩。192 個動作任何時刻大多非法——超界、格子已填。"
         "每次取樣前把非法動作分數設成負無窮，softmax 後機率為零、熵也為零。沒這步，"
         "策略會把大半能力浪費在學不要下非法步。同一個遮罩兼任遊戲結束偵測：沒有合法步"
         "就結束。")

# ===== Slide 10 — My PPO setup ==============================================
s = new_slide()
content_bg(s)
title_band(s, "PPO · IMPLEMENTATION 2", "My PPO Setup — Custom, not SB3",
           "10 / 14")
# top banner
rect(s, 60, 150, PAGE_W - 120, 118, fill=NAVY, rounded=True, radius=0.03)
txt(s, 92, 168, PAGE_W - 184, 20,
    [{"t": "ALGORITHM · CUSTOM PPO", "size": 11, "bold": True, "color": TEAL_LT,
      "spacing": 1.3}])
txt(s, 92, 194, PAGE_W - 184, 66,
    [{"runs": [{"t": "Written from scratch so it shares the EXACT CNN backbone "
                "with the DQN member ", "size": 15.5, "color": WHITE},
               {"t": "— SB3 won't allow that cleanly. ", "size": 15.5,
                "color": WHITE},
               {"t": "Every hyperparameter still matches SB3 MaskablePPO "
                "defaults, so neither method was hand-tuned to win.",
                "size": 15.5, "bold": True, "color": AMBER}], "line": 1.3}])
params = [
    ("clip_range = 0.2", "ratio clipped to [0.8, 1.2]", TEAL),
    ("γ = 0.99 · λ = 0.95", "discount + GAE smoothing", TEAL),
    ("lr = 3e-4 → 0", "Adam, linear decay", AMBER),
    ("8 envs × 128 = 1024", "transitions per iteration", TEAL),
    ("10 epochs × 64", "minibatch reuse; adv normalized per minibatch", TEAL),
    ("health metrics", "approx-KL · clip fraction · explained variance", AMBER),
]
cw, ch, gx, gy, x0, y0 = 388, 96, 26, 16, 60, 288
for i, (h, b, ac) in enumerate(params):
    cx = x0 + (i % 3) * (cw + gx)
    cyy = y0 + (i // 3) * (ch + gy)
    rect(s, cx, cyy, cw, ch, fill=CARD, shadow=True)
    left_accent(s, cx, cyy, ch, ac)
    txt(s, cx + 24, cyy + 18, cw - 44, 26,
        [{"t": h, "size": 16, "bold": True, "color": TEAL_DK, "font": MONO}])
    txt(s, cx + 24, cyy + 50, cw - 44, 38,
        [{"t": b, "size": 12.5, "color": INK, "line": 1.2}])
notes(s, "一個值得提的設計決定：我從零寫 PPO，而不是用 SB3。原因是公平——專案要比較 "
         "PPO 與 DQN，只有共用完全相同的網路主幹才有意義，SB3 沒辦法乾淨做到。為了誠實，"
         "每個超參數都對齊 SB3 預設值，兩邊都沒為了贏而調。每輪在 8 個平行遊戲收集 1024 "
         "步、重複利用十個 epoch。每次更新記錄 approx-KL、clip fraction、explained "
         "variance 確認穩定。")

# ===== Slide 11 — Shared backbone ===========================================
s = new_slide()
content_bg(s)
title_band(s, "PPO · IMPLEMENTATION 3",
           "The Shared CNN Backbone", "11 / 14")
rect(s, 60, 150, PAGE_W - 120, 92, fill=NAVY, rounded=True, radius=0.03)
txt(s, 92, 166, PAGE_W - 184, 18,
    [{"t": "PURPOSE", "size": 11, "bold": True, "color": TEAL_LT,
      "spacing": 1.4}])
txt(s, 92, 188, PAGE_W - 184, 46,
    [{"t": "The base actor–critic — the SAME backbone the DQN member uses. Any "
           "PPO-vs-DQN gap is attributable to the algorithm, not the encoder.",
      "size": 15, "color": WHITE, "line": 1.25}])
cols = [
    ("01", "Board branch", "board (1,8,8)\n→ 2× Conv 3×3\n→ (64,8,8)", TEAL),
    ("02", "Pieces branch", "pieces (3,5,5)\n→ weight-shared encoder\n"
     "→ pad → (48,8,8)", TEAL),
    ("03", "Fuse + heads", "fusion conv → FC(4096→128)\n+ pieces_left(3) → "
     "FC(131→128)\n→ actor (192) + critic (1)", AMBER),
]
cw, gx, x0, cy = 388, 26, 60, 262
for i, (n, h, b, ac) in enumerate(cols):
    cx = x0 + i * (cw + gx)
    rect(s, cx, cy, cw, 56, fill=ac)
    txt(s, cx + 22, cy, 70, 56,
        [{"t": n, "size": 20, "bold": True, "color": WHITE}],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, cx + 78, cy, cw - 96, 56,
        [{"t": h, "size": 15, "bold": True, "color": WHITE}],
        anchor=MSO_ANCHOR.MIDDLE)
    rect(s, cx, cy + 56, cw, 196, fill=CARD, shadow=True)
    txt(s, cx + 26, cy + 84, cw - 48, 150,
        [{"t": b, "size": 14.5, "color": INK, "font": MONO, "line": 1.45}])
by = 540
rect(s, 60, by, PAGE_W - 120, 84, fill=NAVY, rounded=True, radius=0.04)
txt(s, 92, by + 16, PAGE_W - 184, 18,
    [{"t": "SPATIAL FUSION", "size": 11, "bold": True, "color": AMBER,
      "spacing": 1.4}])
txt(s, 92, by + 38, PAGE_W - 184, 42,
    [{"t": "Fusion is performed while features retain their 2-D topology "
           "(8×8) — the network lines up empty cells with piece shapes "
           "positionally, before any flattening.", "size": 14,
      "color": WHITE, "line": 1.25}])
notes(s, "這是網路——基礎 actor-critic，也是 DQN 用的同一主幹。棋盤過兩層卷積；三塊"
         "棋子過權重共享 encoder、補零到相同空間大小，融合卷積結合，讓網路在空間上推理"
         "每塊棋子能放哪。接 pieces_left，再分 actor 與 critic 兩個 head。編碼器兩邊"
         "相同，所以 PPO 與 DQN 的差距來自演算法、非網路。融合是在還保持二維時做的——"
         "攤平前就把空格和棋子形狀在位置上對齊。")

# ===== Slide 12 — Variant 1 heuristics ======================================
s = new_slide()
content_bg(s)
title_band(s, "PPO · VARIANT 1", "Heuristic-Aware Actor–Critic", "12 / 14")
hero_card(s, 60, 150, 560, 470,
          "VARIANT 1", "+ 40-d hand-\ncrafted features",
          sub="hand the agent the answer, don't make it relearn the board",
          accent=TEAL,
          body=[{"t": "Added alongside the shared backbone without modifying "
                      "it — the PPO-vs-DQN comparison stays fair.",
                 "size": 14, "color": INK_ON_DARK, "line": 1.35}])
cards = [
    ("The 40 features", "column heights · holes · row/col fill · bumpiness · "
     "# legal moves — all emitted by the env, normalized to [0,1].", TEAL),
    ("How they merge", "small MLP (40→32), concatenated into the shared "
     "layer:  fusion(128) + pieces_left(3) + heuristics(32) → FC(163→128).",
     TEAL),
    ("Why it helps", "the agent no longer has to learn \"this column is too "
     "tall\" from pixels — it gets a head start.", AMBER),
]
rx, ry = 652, 150
for h, b, ac in cards:
    hh = 150
    rect(s, rx, ry, 621, hh, fill=CARD, shadow=True)
    left_accent(s, rx, ry, hh, ac)
    txt(s, rx + 28, ry + 22, 565, 24,
        [{"t": h, "size": 16, "bold": True, "color": INK}])
    txt(s, rx + 28, ry + 54, 568, 88,
        [{"t": b, "size": 13.5, "color": INK, "line": 1.3}])
    ry += hh + 10
notes(s, "第一個變體加入手工特徵。環境本來就算出 40 個有用數字——各欄高度、洞、表面"
         "崎嶇度、剩幾個合法步。與其逼 CNN 從原始輸入重新發現，我透過小 MLP 餵進去、"
         "合併到共享層，給 agent 起跑優勢。關鍵是在不修改共享主幹的前提下額外加上，所以"
         "比較仍公平。")

# ===== Slide 13 — Variant 2 afterstate ======================================
s = new_slide()
content_bg(s)
title_band(s, "PPO · VARIANT 2", "Afterstate Evaluation", "13 / 14")
hero_card(s, 60, 150, 560, 470,
          "VARIANT 2 · ~2.7k PARAMS", "Score the board\nafter the move",
          sub="simulate each placement, score the resulting board",
          accent=AMBER,
          body=[{"t": "A classic idea for Tetris-like games (deterministic "
                      "placement), adapted to Block Blast's 192-action space.",
                 "size": 14, "color": INK_ON_DARK, "line": 1.35}])
cards = [
    ("9 afterstate features", "per action: lines cleared, eroded cells, "
     "holes, max/mean height, bumpiness, row/col transitions, near-full lines.",
     TEAL),
    ("Tiny actor + critic", "actor = one shared Linear(9→1) scoring all 192 "
     "actions; critic = 40-d heuristics → small MLP → V(s).", TEAL),
    ("Raw integer scale on purpose", "normalizing to [0,1] shrinks logit gaps "
     "below the entropy bonus → policy freezes uniform — never clears the "
     "entropy floor.", AMBER),
]
rx, ry = 652, 150
for h, b, ac in cards:
    hh = 150
    rect(s, rx, ry, 621, hh, fill=CARD, shadow=True)
    left_accent(s, rx, ry, hh, ac)
    txt(s, rx + 28, ry + 20, 565, 24,
        [{"t": h, "size": 16, "bold": True, "color": INK}])
    txt(s, rx + 28, ry + 52, 568, 92,
        [{"t": b, "size": 13.5, "color": INK, "line": 1.28}])
    ry += hh + 10
notes(s, "第二個變體最有趣，而且非常小——不到三千個參數。想法是別評估現在的狀態，改成"
         "往前看一步：對每個合法動作，模擬放置，再用九個簡單特徵描述結果棋盤。一個共享"
         "線性層為全部 192 個動作打分。微妙但重要：我刻意不正規化這些特徵。actor 只是"
         "小初始化的線性層，若壓到 0 到 1，動作分數一開始幾乎一樣，差距小於熵獎勵，策略"
         "凍結在均勻隨機、跨不過熵地板。保留原始整數（消兩行就是 2，不是 0.13）才能從"
         "第一次更新就分辨好步壞步。afterstate 評估是 Tetris 類遊戲的經典思路。")

# ===== Slide 14 — Recap =====================================================
s = new_slide()
content_bg(s)
title_band(s, "PPO · RECAP", "Recap", "14 / 14")
cols = [
    ("What it is", ["model-free", "policy-based (actor–critic)", "on-policy",
                    "A2C family — but safer & more efficient"], TEAL),
    ("The math · 4 pieces", ["ratio  (importance sampling)",
                             "advantage  (GAE)",
                             "clipped surrogate objective",
                             "3-term loss"], TEAL),
    ("Implementation", ["action masking (mandatory)",
                        "from-scratch PPO, SB3-matched",
                        "shared CNN backbone",
                        "+ 2 variants: heuristic · afterstate"], AMBER),
]
cw, gx, x0, cy = 388, 26, 60, 158
for i, (h, items, ac) in enumerate(cols):
    cx = x0 + i * (cw + gx)
    rect(s, cx, cy, cw, 300, fill=CARD, shadow=True)
    rect(s, cx, cy, cw, 8, fill=ac)
    txt(s, cx + 26, cy + 28, cw - 50, 30,
        [{"t": h, "size": 18, "bold": True, "color": INK}])
    rect(s, cx + 26, cy + 64, 42, 4, fill=ac)
    paras = [{"runs": [{"t": "•  ", "size": 14, "color": ac},
                       {"t": it, "size": 14, "color": INK}],
              "line": 1.2, "space_after": 12} for it in items]
    txt(s, cx + 26, cy + 84, cw - 50, 200, paras)
by = 478
rect(s, 60, by, PAGE_W - 120, 142, fill=NAVY, rounded=True, radius=0.03)
txt(s, 92, by + 24, PAGE_W - 184, 20,
    [{"t": "TWO IDEAS MAKE IT WORK", "size": 11, "bold": True, "color": AMBER,
      "spacing": 1.4}])
txt(s, 92, by + 50, PAGE_W - 184, 50,
    [{"runs": [{"t": "Importance sampling", "size": 16, "bold": True,
                "color": WHITE},
               {"t": "  → reuse each batch        ", "size": 16,
                "color": INK_ON_DARK},
               {"t": "Clipping", "size": 16, "bold": True, "color": WHITE},
               {"t": "  → keep those updates stable", "size": 16,
                "color": INK_ON_DARK}], "line": 1.3}])
txt(s, 92, by + 100, PAGE_W - 184, 28,
    [{"t": "Results compared in the next section.", "size": 13.5,
      "italic": True, "color": TEAL_LT}])
notes(s, "總結：PPO 是 model-free、policy-based、on-policy 的 actor-critic——和 A2C "
         "同類，但更安全、更省樣本，靠兩個想法：重要性採樣讓我們重複利用每批資料，clip "
         "讓更新穩定。數學四塊：ratio、advantage、裁剪目標、三項損失。我從零實作以共用"
         "主幹，對齊 SB3 讓超參數誠實，做了兩個變體——手工特徵、後繼狀態評估。下一章節"
         "比較實際表現。")

# ----------------------------------------------------------------------------
out_pptx = os.path.join(HERE, "PPO_Section05.pptx")
prs.save(out_pptx)
print("wrote", out_pptx, "—", len(prs.slides._sldIdLst), "slides")

# ----------------------------------------------------------------------------
# Speaker-script Word document (full EN + 中文 narration, one block per slide)
# ----------------------------------------------------------------------------
from docx import Document
from docx.shared import Pt as DPt, RGBColor as DRGB
from docx.oxml.ns import qn as dqn

SCRIPTS = [
    ("0 · Section divider / 章節頁",
     "This section is on PPO — Proximal Policy Optimization, the algorithm I "
     "used for our agent. I'll cover what kind of method it is, the math that "
     "makes it work, and how I implemented it from scratch.",
     "這個章節講 PPO——Proximal Policy Optimization，也就是我為我們 agent 用的"
     "演算法。我會講它是哪一類方法、讓它運作的數學、以及我怎麼從零實作它。"),
    ("1 · What Kind of Algorithm Is PPO?",
     "First, where PPO sits. It's model-free — it learns purely from "
     "experience and never tries to predict the game's dynamics. It's "
     "policy-based, meaning it directly optimizes the policy, the actor, "
     "rather than learning Q-values like DQN does. And it's on-policy — it "
     "trains on data the current policy just collected, then throws it away. "
     "It's in the same actor–critic family as A2C and A3C, but it's noticeably "
     "more stable and squeezes more learning out of each batch. Notice the "
     "contrast with our DQN, which is value-based and off-policy — that "
     "contrast is exactly what our project compares.",
     "首先是 PPO 的定位。它是 model-free——純粹從經驗學，從不試圖預測遊戲的動態。"
     "它是 policy-based，直接優化策略（actor），而不是像 DQN 那樣學 Q 值。它是 "
     "on-policy——用當前策略剛收集的資料訓練，然後丟掉。它和 A2C、A3C 同屬 "
     "actor-critic 家族，但明顯更穩定、能從每批資料榨出更多學習。注意它和 DQN 的"
     "對比，DQN 是 value-based、off-policy——這個對比正是我們專案要比較的。"),
    ("2 · The PPO Training Pipeline",
     "Before the math, let me show you how PPO trains — it's a loop with three "
     "stages. Stage one, rollout: the current actor plays N steps and we store "
     "states, actions, the old log-probabilities, rewards, values, and done "
     "flags. Stage two, GAE: with the network frozen, a backward pass gives "
     "advantages and returns. Stage three, optimize: we shuffle that batch and "
     "train on it. Coming from A2C you'll notice stage three reuses the same "
     "batch for K epochs instead of throwing it away after one update. That's "
     "a big efficiency win, but it should make you nervous: reusing on-policy "
     "data is exactly what breaks vanilla policy gradient. So how does PPO get "
     "away with it? Two tricks — importance sampling and clipping — and that's "
     "precisely what the next few slides, the math, are about.",
     "進數學前，先給你看 PPO 怎麼訓練——三階段迴圈。階段一 rollout：當前 actor 跑 "
     "N 步，把狀態、動作、舊 log 機率、獎勵、價值、done 都存下來。階段二 GAE：凍結"
     "網路，一次反向掃描算出優勢與回報。階段三 optimize：打散那批資料、拿來訓練。"
     "從 A2C 過來你會注意到，第三階段把同一批資料重複利用 K 個 epoch，而不是更新一"
     "次就丟。效率大增，但它應該讓你緊張：重複利用 on-policy 資料正是搞垮傳統 policy "
     "gradient 的元兇。那 PPO 怎麼做到？兩招——重要性採樣與 clip——正是接下來數學"
     "那幾頁要講的。"),
    ("3 · The Math at a Glance",
     "We just saw the loop and the two tricks that keep it safe — here's all "
     "the math on one slide. PPO is really just four pieces: the ratio that "
     "measures how much the policy changed, an advantage that measures whether "
     "an action was good, the clipped objective that updates the policy "
     "safely, and a total objective that ties in the critic and exploration. "
     "Don't worry about the symbols yet — I'll take them one at a time. Keep "
     "this slide as the map.",
     "我們剛看完迴圈、還有讓它安全的兩招——這頁把整個數學放在一起。PPO 其實就四塊："
     "量化策略改變多少的 ratio、量化動作好不好的 advantage、安全更新策略的裁剪目標、"
     "把 critic 和探索綁進來、要最大化的總目標。符號先別擔心，我會一個一個講。這頁"
     "當作地圖。"),
    ("4 · Probability Ratio = Importance Sampling",
     "Piece one: the ratio — and this is the first of our two tricks, "
     "importance sampling. It's the new policy's probability for an action "
     "divided by the old policy's probability for that same action. One means "
     "nothing changed; above one means the policy now likes that action more. "
     "Because we're reusing data the old policy collected, this ratio is "
     "exactly what corrects for the mismatch — it's what makes the reuse "
     "mathematically valid. We compute it in log-space for numerical "
     "stability. This ratio is also the quantity PPO will clip — the second "
     "trick, coming up.",
     "第一塊：ratio——這也是兩招中的第一招，重要性採樣。它是新策略對某動作的機率，"
     "除以舊策略對同一動作的機率。等於 1 沒變；大於 1 代表更喜歡。因為在重複利用舊"
     "策略收集的資料，這個比率正是修正落差、讓重複利用在數學上成立的東西。在 log "
     "空間算以保數值穩定。這個 ratio 也是接下來要 clip 的量——第二招。"),
    ("5 · Advantage & GAE",
     "Piece two: advantage — was the action better or worse than the critic "
     "expected? Positive, we reinforce it; negative, we discourage it. The "
     "real question is how to estimate it, and there are two extremes. A "
     "one-step TD estimate uses just the next real reward plus the critic's "
     "own value guess for everything after — relying on a guess like that is "
     "called bootstrapping. It's stable, low variance, but biased by that "
     "guess. Monte-Carlo is the opposite: use the real rewards all the way to "
     "the end of the episode, no guessing — unbiased, but high variance. GAE "
     "generalizes both: a weighted blend controlled by lambda — zero is pure "
     "TD, one is pure Monte-Carlo, and we use 0.95. If you unroll the "
     "recursion, the advantage is just an exponentially-weighted sum of all "
     "the future TD errors, with decay gamma-times-lambda — so the further out "
     "a step is, the less its noise counts. That one number is the dial that "
     "trades variance against bias. At an episode boundary the recursion "
     "resets, so value doesn't bleed across the end of a game.",
     "第二塊：advantage——比 critic 預期好還是壞？正的強化、負的抑制。真正的問題是"
     "怎麼估，這有兩個極端。單步 TD 只用下一個真實獎勵加 critic 自己的價值猜測——這種"
     "依賴猜測叫 bootstrap，穩定、低變異但被帶偏。Monte-Carlo 相反：用到 episode "
     "結束的真實獎勵、不猜，無偏但高變異。GAE 把兩者一般化：由 lambda 控制的加權混合"
     "——0 是純 TD、1 是純 Monte-Carlo，我們用 0.95。把遞迴展開，advantage 就是所有"
     "未來 TD error 的指數加權和，衰減 gamma 乘 lambda——越遠雜訊影響越小。這個數字"
     "就是權衡變異與偏差的旋鈕。在 episode 邊界遞迴重置，價值不會跨越遊戲結束外溢。"),
    ("6 · The Clipped Surrogate Objective",
     "Piece three is the heart of PPO, and it's the second trick — the answer "
     "to the puzzle from the pipeline slide. We take two versions of the same "
     "term: the unclipped one, the ordinary policy gradient, and a clipped one "
     "where the ratio is forced to stay between 0.8 and 1.2. We keep the "
     "smaller — the pessimistic — of the two. That lower bound removes any "
     "reward for moving the policy too far in one update, which is exactly "
     "what makes reusing the batch for many epochs safe. This single cheap "
     "operation replaces the heavy trust-region machinery that TRPO needed. "
     "The next slide shows how it brakes the update in each case.",
     "第三塊是 PPO 核心，也是第二招——流程那頁疑問的答案。同一項取兩版本：未裁剪的"
     "（一般 policy gradient），與裁剪的（比率被強制留在 0.8 到 1.2）。取較小、較"
     "悲觀的那個。那個下界讓單次把策略移動太遠得不到任何獎勵——這正是讓重複利用多個 "
     "epoch 安全的原因。這一個便宜操作取代了 TRPO 笨重的 trust-region。下一頁說明它"
     "在各情況怎麼踩煞車。"),
    ("7 · Why Clipping Works — Case Analysis",
     "Let's go case by case. For a good action with positive advantage: while "
     "the ratio is within range, it's the normal gradient and we keep "
     "improving. Once the ratio passes 1.2, the min switches to a constant, "
     "its gradient is zero, and we stop — a ceiling that prevents "
     "over-committing. For a bad action with negative advantage it's "
     "symmetric: below 0.8 the gradient flattens — a floor that prevents "
     "over-punishing. Now the safety case: if the policy moves the wrong way "
     "on a bad action and the ratio shoots up, the unclipped term is even more "
     "negative, so the min keeps it — and we get a big gradient pulling the "
     "policy back. So: a ceiling, a floor, but no free pass for going the "
     "wrong direction.",
     "逐情況看。好動作、advantage 為正：比率在範圍內就是正常梯度，持續改善；一旦超過 "
     "1.2，min 切換成常數、梯度為零，就停——天花板，防過度押注。壞動作 advantage 為負"
     "則對稱：低於 0.8 梯度變平——地板，防過度懲罰。安全情況：壞動作往錯方向、比率暴衝"
     "時，未裁剪項更負，min 保留它，給大梯度把策略拉回。所以：天花板、地板，但往錯方向"
     "不給通行證。"),
    ("8 · The Complete Loss",
     "Piece four ties it together. Three terms: the clipped objective trains "
     "the actor; a mean-squared-error term trains the critic to predict the "
     "GAE returns; and an entropy bonus keeps the policy a little random so it "
     "keeps exploring instead of collapsing too early. One detail that often "
     "draws a question — the value term is subtracted in the paper's objective "
     "but added in the code. That's only because the paper writes an objective "
     "to maximize while PyTorch minimizes, so the code negates the whole thing "
     "and every sign flips. Same equation, opposite convention. The "
     "coefficients, 0.5 and 0.01, are the standard values I kept from SB3.",
     "第四塊把全部綁起來。三項：裁剪目標訓練 actor；均方誤差訓練 critic 預測 GAE "
     "returns；熵獎勵讓策略保持一點隨機、持續探索而不太早崩塌。常被問的細節——價值項"
     "在論文的『目標』裡是減、在程式碼裡卻是加。那只是因為論文寫的是要最大化的目標，"
     "而 PyTorch 是最小化，程式碼把整個取負號、每項符號都翻。同一條式子、相反慣例。"
     "係數 0.5 和 0.01 沿用 SB3。"),
    ("9 · Action Masking",
     "Now the implementation, and the single most important detail is action "
     "masking. We have 192 possible actions but at any moment most are illegal "
     "— out of bounds, or the cell's already filled. Before every sampling "
     "step I set the illegal actions' scores to negative infinity, so after "
     "softmax they get exactly zero probability and contribute zero entropy. "
     "Without this the policy would waste most of its capacity just learning "
     "not to make illegal moves. The same mask doubles as game-over detection: "
     "if nothing is legal, the episode is done.",
     "進到實作，最重要的一個細節就是動作遮罩。192 個可能動作，任何時刻大多非法——超出"
     "邊界、或格子已填。每次取樣前我把非法動作的分數設成負無窮，softmax 後機率剛好為"
     "零、熵貢獻也是零。沒這步，策略會把大半能力浪費在學不要下非法步。同一個遮罩兼任"
     "遊戲結束偵測：沒有任何合法步，episode 就結束。"),
    ("10 · My PPO Setup — Custom, not SB3",
     "One design decision worth calling out: I wrote PPO from scratch instead "
     "of using Stable-Baselines3. The reason is fairness — our project "
     "compares PPO to DQN, and that's only meaningful if both share the "
     "identical network backbone, which SB3 doesn't allow cleanly. To stay "
     "honest I matched every hyperparameter to SB3's defaults, so neither side "
     "was secretly tuned to win. Each iteration collects 1024 steps across 8 "
     "parallel games and reuses them for ten epochs. I also log a few health "
     "metrics every update — approximate KL, clip fraction, explained variance "
     "— to confirm training stays stable.",
     "一個值得提的設計決定：我從零寫 PPO，而不是用 Stable-Baselines3。原因是公平——"
     "專案要比較 PPO 和 DQN，只有共用完全相同的網路主幹才有意義，SB3 沒辦法乾淨做到。"
     "為了誠實，每個超參數都對齊 SB3 預設值，兩邊都沒偷偷為了贏而調。每輪在 8 個平行"
     "遊戲收集 1024 步、重複利用十個 epoch。每次更新還記錄近似 KL、clip fraction、"
     "explained variance，確認訓練穩定。"),
    ("11 · The Shared CNN Backbone",
     "Here's the network — the base actor-critic, and the same backbone the "
     "DQN member uses. The board passes through two convolutions; the three "
     "pieces go through a weight-shared encoder, get padded to the same "
     "spatial size, and a fusion convolution combines them while everything is "
     "still 2-D — so the network lines up empty cells with piece shapes "
     "positionally, before any flattening. We append the pieces-left flags, "
     "then split into the actor and critic heads. Because this encoder is "
     "identical on both sides, any performance gap between PPO and DQN comes "
     "from the algorithm, not the network.",
     "這是網路——基礎 actor-critic，也是 DQN 用的同一主幹。棋盤經過兩層卷積；三塊棋子"
     "經過權重共享 encoder、補零到相同空間大小，一個融合卷積在還是二維時把它們結合——"
     "讓網路在攤平前就把空格和棋子形狀在位置上對齊。接上剩餘棋子旗標，再分成 actor 與 "
     "critic 兩個 head。編碼器兩邊完全相同，所以 PPO 與 DQN 的任何差距都來自演算法、"
     "而非網路。"),
    ("12 · Variant 1 — Heuristic-Aware Actor–Critic",
     "The first variant adds hand-crafted features. The environment already "
     "computes 40 useful numbers — column heights, holes, how bumpy the "
     "surface is, how many legal moves remain. Rather than forcing the CNN to "
     "rediscover these from raw input, I feed them through a small MLP and "
     "merge them into the shared layer, giving the agent a head start. "
     "Crucially I add this alongside the shared backbone without modifying it, "
     "so the PPO-versus-DQN comparison stays fair.",
     "第一個變體加入手工特徵。環境本來就算出 40 個有用數字——各欄高度、洞、表面崎嶇"
     "度、剩幾個合法步。與其逼 CNN 從原始輸入重新發現，我透過小 MLP 餵進去、合併到共享"
     "層，給 agent 起跑優勢。關鍵是在不修改共享主幹的前提下額外加上，所以 PPO 對 DQN "
     "的比較仍公平。"),
    ("13 · Variant 2 — Afterstate Evaluation",
     "The second variant is the most interesting, and it's tiny — under three "
     "thousand parameters. The idea is to stop scoring the current state and "
     "instead look one move ahead: for every legal action, simulate the "
     "placement and describe the resulting board with nine simple features, "
     "like lines cleared and holes created. A single shared linear layer "
     "scores all 192 actions. One subtle but important point: I deliberately "
     "don't normalize these features. The actor is a single linear layer with "
     "small initial weights, so if the features were squeezed into zero-to-one, "
     "every action's score would start almost identical — those tiny gaps fall "
     "below the entropy bonus, which actively rewards a uniform policy, and "
     "training never escapes random play. Keeping the raw integer scale — "
     "clearing two lines really is a 2, not 0.13 — lets it break through that "
     "entropy floor and tell good moves from bad from the very first updates. "
     "Afterstate evaluation is a classic idea for Tetris-like games.",
     "第二個變體最有趣，而且非常小——不到三千個參數。想法是別評估現在的狀態，改成往前"
     "看一步：對每個合法動作模擬放置，再用九個簡單特徵描述結果棋盤，像消了幾行、製造"
     "幾個洞。一個共享線性層為全部 192 個動作打分。微妙但重要：我刻意不正規化這些特徵。"
     "actor 只是一個小初始化的線性層，若把特徵壓到 0 到 1，每個動作的分數一開始幾乎"
     "一樣——這些微小差距小於熵獎勵，而熵獎勵會主動獎勵均勻策略，訓練就永遠跳不出隨機。"
     "保留原始整數尺度——消兩行真的就是 2，不是 0.13——讓它突破熵地板，從第一次更新就"
     "分辨好步壞步。afterstate 評估是 Tetris 類遊戲的經典思路。"),
    ("14 · Recap",
     "To wrap up: PPO is a model-free, policy-based, on-policy actor–critic — "
     "the same family as A2C, but safer and more sample-efficient thanks to "
     "two ideas: importance sampling, which lets us reuse each batch, and "
     "clipping, which keeps those updates stable. The math is four pieces: "
     "ratio, advantage, the clipped objective, and the three-term loss. I "
     "implemented it from scratch to share a backbone with DQN, kept the "
     "hyperparameters honest by matching SB3, and built two variants on top — "
     "one with hand-crafted features, one with afterstate evaluation. The next "
     "section compares how they actually performed.",
     "總結：PPO 是 model-free、policy-based、on-policy 的 actor-critic——和 A2C 同類，"
     "但更安全、更省樣本，靠兩個想法：重要性採樣讓我們重複利用每批資料，clip 讓更新"
     "穩定。數學就四塊：ratio、advantage、裁剪目標、三項損失。我從零實作以共用主幹，"
     "對齊 SB3 讓超參數誠實，並做了兩個變體——手工特徵、後繼狀態評估。下一章節比較它們"
     "實際表現。"),
]


def _set_cjk(run, font="Microsoft JhengHei"):
    run.font.name = font
    rpr = run._element.get_or_add_rPr()
    rfonts = rpr.find(dqn('w:rFonts'))
    if rfonts is None:
        rfonts = rpr.makeelement(dqn('w:rFonts'), {})
        rpr.append(rfonts)
    rfonts.set(dqn('w:eastAsia'), font)


doc = Document()
doc.add_heading("PPO — Section 05 · Speaker Script / 講稿", level=0)
intro = doc.add_paragraph(
    "On-slide text is English; this document holds the full narration in "
    "English and 繁體中文 for each slide. Slide notes in the .pptx hold the "
    "Chinese summary; this is the talk-track.")
intro.runs[0].font.size = DPt(10)
intro.runs[0].font.italic = True

for title, en, zh in SCRIPTS:
    h = doc.add_heading(title, level=1)
    pe = doc.add_paragraph()
    re = pe.add_run("EN  ")
    re.bold = True
    re.font.color.rgb = DRGB(0x17, 0x7C, 0x84)
    pe.add_run(en)
    pz = doc.add_paragraph()
    rz = pz.add_run("中文  ")
    rz.bold = True
    rz.font.color.rgb = DRGB(0xC0, 0x6A, 0x10)
    _set_cjk(rz)
    rzt = pz.add_run(zh)
    _set_cjk(rzt)
    doc.add_paragraph("")

out_docx = os.path.join(HERE, "PPO_speaker_script.docx")
doc.save(out_docx)
print("wrote", out_docx)
