"""Generate Block Blast RL presentation slides."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
BG_WHITE   = RGBColor(0xFA, 0xFA, 0xFC)   # near-white slide background
ACCENT     = RGBColor(0x2E, 0x6D, 0xC8)   # blue for titles / accents
ACCENT2    = RGBColor(0x1A, 0x47, 0x8A)   # darker blue
TEXT_DARK  = RGBColor(0x1E, 0x1E, 0x2E)   # near-black body text
DIVIDER    = RGBColor(0xCC, 0xD9, 0xF0)   # light blue divider
PLACEHOLDER_BG = RGBColor(0xE8, 0xEF, 0xFA)  # chart placeholder fill
PLACEHOLDER_BD = RGBColor(0xA0, 0xB8, 0xE0)  # chart placeholder border
TAG_BG     = RGBColor(0xE2, 0xEC, 0xF9)   # tag/chip background

# ── Slide dimensions (Widescreen 16:9) ───────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK_LAYOUT = prs.slide_layouts[6]   # truly blank

# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width if line_width > Pt(0) else Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=Pt(22), bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_bullet_box(slide, bullets, l, t, w, h,
                   size=Pt(20), color=TEXT_DARK, spacing=1.2):
    """Add a textbox with multiple bullet lines (• prefix)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size = size
        run.font.color.rgb = color
    return tb


def set_slide_bg(slide, color=BG_WHITE):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def title_bar(slide, title_text, subtitle=None):
    """Blue top bar with white title text."""
    bar_h = Inches(1.1)
    add_rect(slide, 0, 0, W, bar_h, fill=ACCENT)
    add_text(slide, title_text,
             Inches(0.45), Inches(0.12), Inches(11.5), Inches(0.72),
             size=Pt(34), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.45), Inches(0.75), Inches(11.5), Inches(0.4),
                 size=Pt(18), color=RGBColor(0xCC, 0xDD, 0xFF))
    # thin bottom accent line
    add_rect(slide, 0, bar_h, W, Pt(3), fill=ACCENT2)


def section_tag(slide, label, l, t):
    """Small colored label chip."""
    w, h = Inches(1.6), Inches(0.32)
    add_rect(slide, l, t, w, h, fill=TAG_BG, line_color=ACCENT, line_width=Pt(1))
    add_text(slide, label, l + Inches(0.08), t + Inches(0.02), w, h,
             size=Pt(13), bold=True, color=ACCENT)


def placeholder_chart(slide, l, t, w, h, label="[ Chart Placeholder ]"):
    """Dashed-border box for a chart to be inserted later."""
    box = add_rect(slide, l, t, w, h,
                   fill=PLACEHOLDER_BG,
                   line_color=PLACEHOLDER_BD, line_width=Pt(1.5))
    add_text(slide, label,
             l, t + h / 2 - Inches(0.25), w, Inches(0.5),
             size=Pt(18), italic=True,
             color=RGBColor(0x70, 0x90, 0xC0),
             align=PP_ALIGN.CENTER)
    return box


# ═════════════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(s)

# Full-width top accent
add_rect(s, 0, 0, W, Inches(0.18), fill=ACCENT)
add_rect(s, 0, Inches(0.18), W, Pt(3), fill=ACCENT2)

# Centre card
card_t = Inches(1.4)
card_h = Inches(4.2)
add_rect(s, Inches(1.0), card_t, Inches(11.33), card_h,
         fill=RGBColor(0xFF, 0xFF, 0xFF))

add_text(s, "Block Blast with Reinforcement Learning",
         Inches(1.3), Inches(1.6), Inches(10.7), Inches(1.0),
         size=Pt(38), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_rect(s, Inches(3.5), Inches(2.85), Inches(6.33), Pt(3), fill=ACCENT)

add_text(s, "PPO with Afterstate: A Board-Aware Policy Gradient Approach",
         Inches(1.3), Inches(2.95), Inches(10.7), Inches(0.7),
         size=Pt(22), color=TEXT_DARK, align=PP_ALIGN.CENTER)

add_text(s, "Member A · Member B · Member C · Member D · Member E",
         Inches(1.3), Inches(3.75), Inches(10.7), Inches(0.45),
         size=Pt(18), color=RGBColor(0x55, 0x66, 0x88), align=PP_ALIGN.CENTER)

add_text(s, "2026",
         Inches(1.3), Inches(4.3), Inches(10.7), Inches(0.4),
         size=Pt(16), color=RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.CENTER)

# Bottom strip
add_rect(s, 0, Inches(6.9), W, Inches(0.6), fill=ACCENT)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 2 — Agenda
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(s)
title_bar(s, "Agenda")

items = [
    ("01", "Background", "Game overview & motivation", "2–3 min"),
    ("02", "Environment", "Observation · Action encoding · Masking", "1–2 min"),
    ("03", "PPO-AF Method", "Afterstate concept · Domain motivation · Architecture", "6–7 min"),
    ("04", "Comparison", "Baseline / DQN / PPO / PPO-AF results", "4–5 min"),
    ("05", "Demo", "Live gameplay visualization", "2–3 min"),
    ("06", "Conclusion", "Summary & future work", "1–2 min"),
]

col_x = [Inches(0.45), Inches(1.15), Inches(3.8), Inches(11.5)]
row_h = Inches(0.78)
row_top = Inches(1.35)

for i, (num, title, desc, dur) in enumerate(items):
    t = row_top + i * row_h
    # alternating row shade
    if i % 2 == 0:
        add_rect(s, Inches(0.35), t, Inches(12.63), row_h - Inches(0.05),
                 fill=RGBColor(0xF0, 0xF4, 0xFD))
    # number bubble
    add_rect(s, Inches(0.42), t + Inches(0.14), Inches(0.56), Inches(0.5),
             fill=ACCENT)
    add_text(s, num,
             Inches(0.42), t + Inches(0.14), Inches(0.56), Inches(0.5),
             size=Pt(16), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER)
    # title
    add_text(s, title,
             Inches(1.12), t + Inches(0.08), Inches(2.5), Inches(0.36),
             size=Pt(22), bold=True, color=ACCENT)
    # description
    add_text(s, desc,
             Inches(1.12), t + Inches(0.42), Inches(8.5), Inches(0.3),
             size=Pt(17), color=TEXT_DARK)
    # duration tag
    add_rect(s, Inches(11.35), t + Inches(0.18), Inches(1.4), Inches(0.38),
             fill=TAG_BG, line_color=ACCENT, line_width=Pt(1))
    add_text(s, dur,
             Inches(11.35), t + Inches(0.18), Inches(1.4), Inches(0.38),
             size=Pt(15), color=ACCENT, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 3 — Background
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(s)
title_bar(s, "Background", "Game Overview & Motivation")

# Left column — game rules
add_text(s, "What is Block Blast?",
         Inches(0.45), Inches(1.3), Inches(5.8), Inches(0.4),
         size=Pt(22), bold=True, color=ACCENT)
add_rect(s, Inches(0.45), Inches(1.72), Inches(5.7), Pt(2), fill=ACCENT)

rules = [
    "• 8 × 8 grid puzzle game",
    "• 3 pieces are dealt per round",
    "• Place all 3 pieces on the board",
    "• Complete rows/columns → cleared & scored",
    "• Game ends when no legal placement exists",
]
add_bullet_box(s, rules,
               Inches(0.45), Inches(1.85), Inches(5.7), Inches(2.8),
               size=Pt(20))

# Right column — motivation
add_text(s, "Why Reinforcement Learning?",
         Inches(7.0), Inches(1.3), Inches(5.9), Inches(0.4),
         size=Pt(22), bold=True, color=ACCENT)
add_rect(s, Inches(7.0), Inches(1.72), Inches(5.8), Pt(2), fill=ACCENT)

motiv = [
    "• Large discrete action space (192 actions)",
    "• Requires long-horizon planning",
    "• No analytical solution — learned policy needed",
    "• Ideal testbed for comparing DQN vs. PPO variants",
]
add_bullet_box(s, motiv,
               Inches(7.0), Inches(1.85), Inches(5.9), Inches(2.5),
               size=Pt(20))

# Research question box
add_rect(s, Inches(0.45), Inches(5.15), Inches(12.43), Inches(1.6),
         fill=RGBColor(0xE6, 0xEE, 0xFA))
add_text(s, "Research Question",
         Inches(0.6), Inches(5.25), Inches(4), Inches(0.36),
         size=Pt(18), bold=True, color=ACCENT)
add_text(s,
         "Can a board-aware afterstate representation help PPO learn more efficiently "
         "and achieve higher scores than standard PPO or DQN on Block Blast?",
         Inches(0.6), Inches(5.62), Inches(12.1), Inches(0.9),
         size=Pt(19), color=TEXT_DARK)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 4 — Environment
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(s)
title_bar(s, "Environment Design")

col_w = Inches(3.9)
col_gap = Inches(0.25)
cols_l = [Inches(0.35), Inches(0.35) + col_w + col_gap,
          Inches(0.35) + 2 * (col_w + col_gap)]

headings = ["Observation Space", "Action Encoding", "Action Masking"]
descs = [
    [
        "Dict observation:",
        "  board: (8, 8) float32",
        "  pieces: (3, 5, 5) float32",
        "",
        "Each piece rendered into",
        "a 5 × 5 binary grid.",
        "Empty slots → all zeros.",
    ],
    [
        "Flat integer encoding:",
        "",
        "  action = piece_idx × 64",
        "           + row × 8 + col",
        "",
        "Total: 192 discrete actions",
        "(3 pieces × 8 rows × 8 cols)",
    ],
    [
        "Most actions are illegal:",
        "  out-of-bounds placement",
        "  overlapping filled cells",
        "  already-used piece slot",
        "",
        "env.action_masks() → (192,) bool",
        "Unmasked agents learn garbage.",
    ],
]

for i, (lx, hd, desc) in enumerate(zip(cols_l, headings, descs)):
    # card background
    add_rect(s, lx, Inches(1.25), col_w, Inches(5.6),
             fill=RGBColor(0xFF, 0xFF, 0xFF))
    # coloured top strip
    add_rect(s, lx, Inches(1.25), col_w, Inches(0.46), fill=ACCENT)
    add_text(s, hd, lx + Inches(0.15), Inches(1.28), col_w, Inches(0.42),
             size=Pt(20), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    # content
    add_bullet_box(s, desc,
                   lx + Inches(0.18), Inches(1.85), col_w - Inches(0.3), Inches(4.8),
                   size=Pt(18), color=TEXT_DARK)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 5 — What is Afterstate?
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(s)
title_bar(s, "PPO-AF: What is an Afterstate?", "Concept")

section_tag(s, "Concept", Inches(0.45), Inches(1.25))

add_text(s,
         "An afterstate is the board configuration that results immediately after an action is taken — "
         "before the environment deals new pieces.",
         Inches(0.45), Inches(1.65), Inches(12.43), Inches(0.8),
         size=Pt(21), color=TEXT_DARK)

placeholder_chart(s,
                  Inches(0.45), Inches(2.6), Inches(12.43), Inches(3.6),
                  "[ Diagram Placeholder — draw your own flowchart here ]")

add_text(s,
         "Key insight: afterstates from different (s, a) pairs may be identical, "
         "allowing the critic to generalize across action paths that lead to the same board.",
         Inches(0.45), Inches(6.35), Inches(12.43), Inches(0.8),
         size=Pt(19), italic=True, color=RGBColor(0x33, 0x55, 0x99))

# ═════════════════════════════════════════════════════════════════════════════
# Slide 6 — Why Afterstate for Block Blast?
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(s)
title_bar(s, "PPO-AF: Why Afterstate for Block Blast?", "Domain Motivation")

section_tag(s, "Motivation", Inches(0.45), Inches(1.25))

points = [
    (
        "Deterministic placement",
        "Given (s, a), the resulting board is fully determined.\n"
        "No stochasticity between action and afterstate — a perfect fit.",
    ),
    (
        "Pieces are noise for the critic",
        "The 3 incoming pieces vary randomly each round.\n"
        "A critic conditioned on the board state ignores this irrelevant noise.",
    ),
    (
        "Board captures all long-term value",
        "Future scores depend entirely on the board layout, not on\n"
        "which pieces happened to be present when the action was taken.",
    ),
]

for i, (head, body) in enumerate(points):
    t = Inches(1.75) + i * Inches(1.75)
    add_rect(s, Inches(0.45), t, Inches(0.56), Inches(0.56),
             fill=ACCENT)
    add_text(s, str(i + 1),
             Inches(0.45), t, Inches(0.56), Inches(0.56),
             size=Pt(24), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER)
    add_text(s, head,
             Inches(1.15), t, Inches(11.7), Inches(0.42),
             size=Pt(22), bold=True, color=ACCENT)
    add_text(s, body,
             Inches(1.15), t + Inches(0.44), Inches(11.7), Inches(1.1),
             size=Pt(19), color=TEXT_DARK)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 7 — Architecture
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(s)
title_bar(s, "PPO-AF: Architecture", "Actor → Afterstate Features  ·  Critic → Board State")

section_tag(s, "Architecture", Inches(0.45), Inches(1.25))

# Standard PPO column
add_rect(s, Inches(0.4), Inches(1.7), Inches(5.6), Inches(5.0),
         fill=RGBColor(0xF5, 0xF7, 0xFF))
add_text(s, "Standard PPO",
         Inches(0.4), Inches(1.75), Inches(5.6), Inches(0.4),
         size=Pt(20), bold=True, color=RGBColor(0x77, 0x88, 0xAA),
         align=PP_ALIGN.CENTER)

std_boxes = [
    (Inches(1.0), Inches(2.3),  "Observation (s)", RGBColor(0xDD, 0xE8, 0xF8)),
    (Inches(1.0), Inches(3.15), "Shared Encoder",  RGBColor(0xDD, 0xE8, 0xF8)),
    (Inches(0.55), Inches(4.0),  "Actor head",      RGBColor(0xDD, 0xE8, 0xF8)),
    (Inches(2.8),  Inches(4.0),  "Critic head",     RGBColor(0xDD, 0xE8, 0xF8)),
]
for (lx, ty, label, fc) in std_boxes:
    add_rect(s, lx, ty, Inches(2.1), Inches(0.55),
             fill=fc, line_color=RGBColor(0xAA, 0xBB, 0xDD), line_width=Pt(1))
    add_text(s, label, lx, ty + Inches(0.05), Inches(2.1), Inches(0.5),
             size=Pt(17), color=TEXT_DARK, align=PP_ALIGN.CENTER)
# arrows
for ty in [Inches(2.85), Inches(3.7)]:
    add_text(s, "↓", Inches(1.95), ty, Inches(0.5), Inches(0.35),
             size=Pt(22), color=RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.CENTER)
add_text(s, "↙   ↘", Inches(1.3), Inches(3.7), Inches(2.4), Inches(0.35),
         size=Pt(22), color=RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.CENTER)
add_text(s, "V(s)", Inches(3.15), Inches(4.65), Inches(1.5), Inches(0.35),
         size=Pt(17), italic=True, color=RGBColor(0x66, 0x77, 0x99), align=PP_ALIGN.CENTER)
add_text(s, "π(a|s)", Inches(0.6), Inches(4.65), Inches(2.0), Inches(0.35),
         size=Pt(17), italic=True, color=RGBColor(0x66, 0x77, 0x99), align=PP_ALIGN.CENTER)

# PPO-AF column
add_rect(s, Inches(7.2), Inches(1.7), Inches(5.75), Inches(5.0),
         fill=RGBColor(0xEA, 0xF1, 0xFF))
add_text(s, "PPO-AF (Ours)",
         Inches(7.2), Inches(1.75), Inches(5.75), Inches(0.4),
         size=Pt(20), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

af_boxes = [
    (Inches(7.7),  Inches(2.3),  "Observation (s)",         RGBColor(0xD4, 0xE6, 0xFF)),
    (Inches(7.7),  Inches(3.1),  "Actor Encoder",           ACCENT),
    (Inches(7.7),  Inches(3.9),  "Afterstate Features (ŝ)", RGBColor(0xD4, 0xE6, 0xFF)),
    (Inches(7.7),  Inches(4.7),  "Critic (board-only)",     ACCENT),
]
for j, (lx, ty, label, fc) in enumerate(af_boxes):
    tc = RGBColor(0xFF, 0xFF, 0xFF) if fc == ACCENT else TEXT_DARK
    add_rect(s, lx, ty, Inches(4.7), Inches(0.55), fill=fc)
    add_text(s, label, lx, ty + Inches(0.05), Inches(4.7), Inches(0.5),
             size=Pt(17), bold=(fc == ACCENT), color=tc, align=PP_ALIGN.CENTER)
    if j < 3:
        add_text(s, "↓", Inches(9.9), ty + Inches(0.55), Inches(0.5), Inches(0.4),
                 size=Pt(22), color=ACCENT, align=PP_ALIGN.CENTER)

add_text(s, "V(ŝ)",
         Inches(9.6), Inches(5.35), Inches(1.5), Inches(0.35),
         size=Pt(17), bold=True, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)

# "vs" divider
add_text(s, "vs.", Inches(6.0), Inches(3.8), Inches(1.2), Inches(0.7),
         size=Pt(28), bold=True, color=RGBColor(0xBB, 0xCC, 0xDD),
         align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# Slide 8 — Method Comparison (chart placeholder)
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(s)
title_bar(s, "Method Comparison", "Baseline · DQN · PPO · PPO-AF")

# Summary table header
cols = ["Agent", "Reward Mode", "Mean Score", "Max Score", "Notes"]
col_widths = [Inches(1.8), Inches(1.8), Inches(1.7), Inches(1.7), Inches(5.28)]
table_l = Inches(0.35)
table_t = Inches(1.3)
row_h_t = Inches(0.38)

x = table_l
for i, (col, cw) in enumerate(zip(cols, col_widths)):
    add_rect(s, x, table_t, cw, row_h_t, fill=ACCENT)
    add_text(s, col, x + Inches(0.05), table_t + Inches(0.05), cw, row_h_t,
             size=Pt(16), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    x += cw

# All values from 100-episode evaluation (seed=42, 2026-05-29)
rows_data = [
    ["Random",  "sparse", "1.17",  "8",   "Uniform-from-legal baseline"],
    ["Greedy",  "sparse", "3.97",  "27",  "One-step lookahead baseline"],
    ["DQN",     "sparse", "2.05",  "10",  "Same CNN backbone as PPO"],
    ["PPO",     "sparse", "2.32",  "9",   "Plateaus at 1 M steps"],
    ["PPO",     "dense",  "4.16",  "15",  "Dense shaping, still improving at 1 M"],
    ["PPO-AF",  "sparse", "34.41", "157", "Board-aware critic"],
    ["PPO-AF",  "dense",  "29.66", "140", "Board-aware critic"],
]
row_fills = [RGBColor(0xF5, 0xF7, 0xFF), RGBColor(0xFF, 0xFF, 0xFF)]
for ri, row in enumerate(rows_data):
    ty = table_t + row_h_t + ri * row_h_t
    x  = table_l
    is_ppoaf = (ri >= 5)
    for ci, (cell, cw) in enumerate(zip(row, col_widths)):
        fc = RGBColor(0xD8, 0xE8, 0xFF) if is_ppoaf else row_fills[ri % 2]
        txt_bold = is_ppoaf
        add_rect(s, x, ty, cw, row_h_t, fill=fc)
        add_text(s, cell, x + Inches(0.05), ty + Inches(0.03),
                 cw - Inches(0.05), row_h_t,
                 size=Pt(15), bold=txt_bold, color=TEXT_DARK)
        x += cw

# Two chart placeholders side by side
placeholder_chart(s,
                  Inches(0.35), Inches(4.45), Inches(6.1), Inches(2.7),
                  "[ Training Curve — Score vs. Steps ]")
placeholder_chart(s,
                  Inches(6.6), Inches(4.45), Inches(6.1), Inches(2.7),
                  "[ Bar Chart — Final Score Comparison ]")

# ═════════════════════════════════════════════════════════════════════════════
# Slide 9 — Demo
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(s)
title_bar(s, "Demo", "Live Gameplay Visualization")

add_text(s, "PPO-AF agent playing Block Blast",
         Inches(0.45), Inches(1.3), Inches(12.43), Inches(0.45),
         size=Pt(22), bold=True, color=ACCENT)

placeholder_chart(s,
                  Inches(0.45), Inches(1.9), Inches(12.43), Inches(4.4),
                  "[ Video / GIF — Demo Footage ]")

add_text(s,
         "Run:  uv run python run_experiments.py demo  "
         "(uses latest PPO-AF checkpoint)",
         Inches(0.45), Inches(6.45), Inches(12.43), Inches(0.45),
         size=Pt(17), italic=True, color=RGBColor(0x55, 0x66, 0x99))

# ═════════════════════════════════════════════════════════════════════════════
# Slide 10 — Conclusion
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(s)
title_bar(s, "Conclusion")

# Left: summary
add_text(s, "Summary",
         Inches(0.45), Inches(1.3), Inches(6.0), Inches(0.4),
         size=Pt(22), bold=True, color=ACCENT)
add_rect(s, Inches(0.45), Inches(1.72), Inches(5.9), Pt(2), fill=ACCENT)
summary = [
    "• Block Blast is a challenging discrete-action RL benchmark",
    "• Action masking is critical for all agents",
    "• PPO-AF separates actor (obs) and critic (afterstate)",
    "• Board-conditioned value function reduces critic variance",
    "• PPO-AF achieves higher mean score than standard PPO",
]
add_bullet_box(s, summary,
               Inches(0.45), Inches(1.85), Inches(5.9), Inches(3.2),
               size=Pt(19))

# Divider
add_rect(s, Inches(6.6), Inches(1.3), Pt(2), Inches(4.5), fill=DIVIDER)

# Right: future work
add_text(s, "Future Work",
         Inches(6.9), Inches(1.3), Inches(5.9), Inches(0.4),
         size=Pt(22), bold=True, color=ACCENT)
add_rect(s, Inches(6.9), Inches(1.72), Inches(5.9), Pt(2), fill=ACCENT)
future = [
    "• Behavioural cloning warm-start",
    "• Larger / deeper CNN backbone",
    "• Handcrafted board features (holes, bumpiness)",
    "• Multi-step lookahead planning",
]
add_bullet_box(s, future,
               Inches(6.9), Inches(1.85), Inches(5.9), Inches(2.5),
               size=Pt(19))

# Thank you box
add_rect(s, Inches(0.45), Inches(5.7), Inches(12.43), Inches(1.35),
         fill=ACCENT, line_color=ACCENT2, line_width=Pt(1))
add_text(s, "Thank You  ·  Q & A",
         Inches(0.45), Inches(5.85), Inches(12.43), Inches(0.65),
         size=Pt(32), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
         align=PP_ALIGN.CENTER)
add_text(s, "Member A · Member B · Member C · Member D · Member E",
         Inches(0.45), Inches(6.35), Inches(12.43), Inches(0.4),
         size=Pt(17), color=RGBColor(0xCC, 0xDD, 0xFF),
         align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
out = r"d:\BlockBlastWithRL\BlockBlast_RL_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
